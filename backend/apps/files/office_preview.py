"""Bounded, read-only extraction for modern Office Open XML files."""
from io import BytesIO
from itertools import islice
from pathlib import Path
import zipfile


MAX_SOURCE_BYTES = 10 * 1024 * 1024
MAX_EXPANDED_BYTES = 50 * 1024 * 1024
MAX_ARCHIVE_ENTRIES = 2000
MAX_TEXT_CHARS = 100_000
MAX_CELL_CHARS = 1000
MAX_SECTIONS = 100
MAX_ROWS = 200
MAX_COLUMNS = 50
MAX_CELLS = 5000

SUPPORTED_EXTENSIONS = {'.docx', '.xlsx', '.pptx'}


class OfficePreviewError(Exception):
    """A user-facing validation error while extracting an Office preview."""


def _read_bounded(file_field):
    file_field.open('rb')
    try:
        payload = file_field.read(MAX_SOURCE_BYTES + 1)
    finally:
        file_field.close()
    if len(payload) > MAX_SOURCE_BYTES:
        raise OfficePreviewError('Office 文件超过 10 MB，无法在线预览')
    return payload


def _validate_archive(payload):
    try:
        with zipfile.ZipFile(BytesIO(payload)) as archive:
            entries = archive.infolist()
            if len(entries) > MAX_ARCHIVE_ENTRIES:
                raise OfficePreviewError('Office 文件内部条目过多，已停止预览')
            expanded = sum(entry.file_size for entry in entries)
            if expanded > MAX_EXPANDED_BYTES:
                raise OfficePreviewError('Office 文件解压内容过大，已停止预览')
            for entry in entries:
                normalized = Path(entry.filename.replace('\\', '/'))
                if normalized.is_absolute() or '..' in normalized.parts:
                    raise OfficePreviewError('Office 文件包含不安全的内部路径')
    except zipfile.BadZipFile as exc:
        raise OfficePreviewError('Office 文件已损坏或格式不正确') from exc


def _clip(value):
    if value is None:
        return ''
    text = str(value).replace('\x00', '').strip()
    return text[:MAX_CELL_CHARS]


def _budgeted_append(target, value, budget):
    text = _clip(value)
    if not text or budget[0] <= 0:
        return False
    text = text[:budget[0]]
    target.append(text)
    budget[0] -= len(text)
    return budget[0] > 0


def _table_rows(rows, budget, cell_counter):
    output = []
    for row_index, row in enumerate(rows):
        if row_index >= MAX_ROWS or budget[0] <= 0 or cell_counter[0] >= MAX_CELLS:
            break
        values = []
        for column_index, value in enumerate(row):
            if column_index >= MAX_COLUMNS or cell_counter[0] >= MAX_CELLS:
                break
            text = _clip(value)
            text = text[:budget[0]]
            values.append(text)
            budget[0] -= len(text)
            cell_counter[0] += 1
            if budget[0] <= 0:
                break
        if any(values):
            output.append(values)
    return output


def _extract_docx(payload):
    from docx import Document

    document = Document(BytesIO(payload))
    budget = [MAX_TEXT_CHARS]
    paragraphs = []
    for paragraph in document.paragraphs[:1000]:
        if not _budgeted_append(paragraphs, paragraph.text, budget):
            break
    tables = []
    cells = [0]
    for table in document.tables[:20]:
        rows = _table_rows(
            ([cell.text for cell in row.cells] for row in table.rows),
            budget,
            cells,
        )
        if rows:
            tables.append(rows)
        if budget[0] <= 0 or cells[0] >= MAX_CELLS:
            break
    return {
        'type': 'docx',
        'sections': [{'title': '文档内容', 'paragraphs': paragraphs, 'tables': tables}],
        'truncated': budget[0] <= 0 or len(document.paragraphs) > 1000,
    }


def _extract_xlsx(payload):
    from openpyxl import load_workbook

    workbook = load_workbook(
        BytesIO(payload), read_only=True, data_only=True, keep_links=False,
    )
    budget = [MAX_TEXT_CHARS]
    cells = [0]
    sections = []
    try:
        for worksheet in workbook.worksheets[:20]:
            rows = _table_rows(
                (
                    [cell.value for cell in row]
                    for row in worksheet.iter_rows(
                        min_row=1,
                        max_row=min(worksheet.max_row or 1, MAX_ROWS),
                        max_col=min(worksheet.max_column or 1, MAX_COLUMNS),
                    )
                ),
                budget,
                cells,
            )
            sections.append({'title': worksheet.title, 'paragraphs': [], 'tables': [rows] if rows else []})
            if budget[0] <= 0 or cells[0] >= MAX_CELLS:
                break
    finally:
        workbook.close()
    return {
        'type': 'xlsx',
        'sections': sections,
        'truncated': len(workbook.sheetnames) > 20 or budget[0] <= 0 or cells[0] >= MAX_CELLS,
    }


def _extract_pptx(payload):
    from pptx import Presentation

    presentation = Presentation(BytesIO(payload))
    budget = [MAX_TEXT_CHARS]
    cells = [0]
    sections = []
    for index, slide in enumerate(islice(presentation.slides, MAX_SECTIONS), start=1):
        paragraphs = []
        tables = []
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                if not _budgeted_append(paragraphs, shape.text, budget):
                    break
            if getattr(shape, 'has_table', False):
                rows = _table_rows(
                    ([cell.text for cell in row.cells] for row in shape.table.rows),
                    budget,
                    cells,
                )
                if rows:
                    tables.append(rows)
            if budget[0] <= 0:
                break
        sections.append({'title': f'第 {index} 页', 'paragraphs': paragraphs, 'tables': tables})
        if budget[0] <= 0 or cells[0] >= MAX_CELLS:
            break
    return {
        'type': 'pptx',
        'sections': sections,
        'truncated': (
            len(presentation.slides) > MAX_SECTIONS
            or budget[0] <= 0
            or cells[0] >= MAX_CELLS
        ),
    }


def extract_office_preview(file_asset):
    extension = Path(file_asset.name or file_asset.file.name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise OfficePreviewError('仅支持 DOCX、XLSX、PPTX 的只读在线预览')
    if file_asset.size and file_asset.size > MAX_SOURCE_BYTES:
        raise OfficePreviewError('Office 文件超过 10 MB，无法在线预览')

    payload = _read_bounded(file_asset.file)
    _validate_archive(payload)
    extractors = {
        '.docx': _extract_docx,
        '.xlsx': _extract_xlsx,
        '.pptx': _extract_pptx,
    }
    try:
        preview = extractors[extension](payload)
    except OfficePreviewError:
        raise
    except Exception as exc:
        raise OfficePreviewError('Office 文件内容无法解析') from exc
    return {
        **preview,
        'name': file_asset.name,
        'limits': {
            'source_bytes': MAX_SOURCE_BYTES,
            'text_chars': MAX_TEXT_CHARS,
            'rows_per_table': MAX_ROWS,
            'columns_per_table': MAX_COLUMNS,
        },
    }
