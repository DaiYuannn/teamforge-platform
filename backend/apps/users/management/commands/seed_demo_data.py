"""
生成完整演示数据的管理命令

用法:
    python manage.py seed_demo_data          # 生成演示数据（带确认提示）
    python manage.py seed_demo_data --clean  # 先清除所有演示数据再重新生成
    python manage.py seed_demo_data --force  # 跳过确认提示

数据概览:
    - 账号: 1 个操作老师 + 3 个查看老师 + 5 个负责人
      + 2 个平行主贡献者 + 35 个普通成员（另有 1 个技术管理员）
    - 团队: 1 个总团队；项目与比赛参赛条目表达实际组队关系
    - 项目: 7 个（2 个历史退役项目 + 5 个当前项目）
    - 比赛: 2021 年起的逐届赛事与参赛条目
    - 任务: 每项目 5 个（覆盖完整状态分布）
    - 文件: 项目基础文档、每个参赛条目独立的计划书/PPT、授权证书
    - 生命周期: 团队成员、项目成员、项目阶段历史
    - 导入历史: 8 个模块的成功、预览、失败与回滚记录
    - 公开门户: 显式公开决策、精选与成员授权
    - 经费: 每项目 1 条预算 + 2~3 条支出
    - 技能标签: 15 个 + 部分成员技能
    - 灵活工时: ~15 条
    - 贡献记录: 15 条
    - 成员排序: 2 个项目（1 draft / 1 confirmed）+ 2 条异议
    - 定时报表: 3 个计划 + 3 个成功执行文件（XLSX/DOCX/PDF）
    - 成果档案: 40 个专利 + 25 个软著 + 7 个科技查新
    - 敏感资料: 3 条 + 2 条访问申请
    - 通知: 10 条
    - 公告: 3 条
    - 操作日志: 20 条
"""
import calendar
import json
import random
import zipfile
from collections import Counter
from datetime import date, datetime, time, timedelta
from decimal import Decimal
from io import BytesIO
from pathlib import Path

from docx import Document
from django.conf import settings
from django.core.files.base import ContentFile
from django.core.management.base import BaseCommand, CommandError
from django.db import transaction
from django.db.models import Q
from django.utils import timezone
from openpyxl import Workbook
from PIL import Image as PILImage, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

from apps.users.models import User, UserLifecycleEvent, UserPreference
from apps.projects.models import (
    Project,
    ProjectMember,
    ProjectMembershipEvent,
    ProjectStageLog,
)
from apps.competitions.models import (
    Competition,
    CompetitionAward,
    CompetitionEvent,
    CompetitionParticipant,
)
from apps.tasks.models import Task
from apps.finance.models import FinanceBudget, FinanceExpense, FinanceIncome, FinanceReceipt
from apps.files.models import FileAsset, FileVersion
from apps.imports.models import ImportTask
from apps.imports.material_archive import preview_material_archive
from apps.common.team_models import Team, TeamMember, TeamMembershipEvent
from apps.dashboard.portal_models import PortalPublication, PortalSettings
from apps.members.models import SkillTag, MemberSkill, FlexibleWorkSchedule
from apps.contributions.models import (
    Contribution, MemberRanking, RankingObjection,
)
from apps.intellectual_property.models import (
    IntellectualPropertyApplication, IPApplicationContributor,
    IPReturnRecord, IPMaterialVersion, IPObjection,
)
from apps.exports.custom_report_models import CustomReport
from apps.exports.scheduled_report_models import (
    ScheduledReport,
    ScheduledReportExecution,
)
from apps.exports.scheduled_report_service import compute_next_run
from apps.sensitive.models import SensitiveData, SensitiveAccessRequest
from apps.notifications.models import Announcement, Notification
from apps.audit.models import OperationLog


DEMO_PROJECT_PREFIX = 'TEAM-DEMO-'
DEMO_IP_PREFIX = 'IP-TEAM-DEMO-'
DEMO_MARKER = '【团队演示】'
DEMO_IMPORT_DIRNAME = 'seed_demo_data'
DEMO_TEAM_CODE = 'TEAM-DEMO-ORG'
DEMO_SQUAD_CODES = ()
LEGACY_DEMO_SQUAD_CODES = (
    'TEAM-DEMO-SQUAD-PRODUCT',
    'TEAM-DEMO-SQUAD-DATA',
    'TEAM-DEMO-SQUAD-OPERATIONS',
)
DEMO_ACCOUNT_EMAILS = (
    'admin@demo.com',
    'teacher1@demo.com',
    'teacher2@demo.com',
    'teacher3@demo.com',
    'teacher4@demo.com',
    'leader1@demo.com',
    'leader2@demo.com',
    'leader3@demo.com',
    'leader4@demo.com',
    'leader5@demo.com',
    'contributor1@demo.com',
    'contributor2@demo.com',
    *(f'member{index}@demo.com' for index in range(1, 36)),
)
LEGACY_COMPETITION_ACCOUNT_EMAILS = (
    'leader6@demo.com',
)
LEGACY_DEMO_ACCOUNT_EMAILS = (
    'approver@demo.com',
    *(f'member{index}@demo.com' for index in range(36, 53)),
)
LEGACY_COMPETITION_REPORT_NAMES = (
    '每日项目经营概览',
    '每周任务交付报告',
    '每月经费执行报告',
    '每周比赛推进简报',
)

FILE_SPECS = (
    (
        '项目概览',
        '.pdf',
        'application/pdf',
        FileAsset.Level.PUBLIC,
    ),
    (
        '工作计划',
        '.docx',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        FileAsset.Level.INTERNAL,
    ),
    (
        '数据简表',
        '.xlsx',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        FileAsset.Level.INTERNAL,
    ),
    (
        '路演材料',
        '.pptx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation',
        FileAsset.Level.INTERNAL,
    ),
)


# ============================================================================
# 中文姓名 / 数据池
# ============================================================================
SURNAMES = [
    '王', '李', '张', '刘', '陈', '杨', '赵', '黄', '周', '吴', '徐', '孙',
    '胡', '朱', '高', '林', '何', '郭', '马', '罗', '梁', '宋', '郑', '谢',
    '韩', '唐', '冯', '于', '董', '萧', '程', '曹', '袁', '邓', '许', '傅',
    '沈', '曾', '彭', '吕', '苏', '卢', '蒋', '蔡', '贾', '丁', '魏', '薛',
    '叶', '阎', '余', '潘', '杜', '戴', '夏', '钟', '汪', '田', '任', '姜',
    '范', '方', '石', '姚', '谭', '廖', '邹', '熊', '金', '陆', '郝', '孔',
    '白', '崔', '康', '毛', '邱', '秦', '江', '史', '顾', '侯', '邵', '孟',
    '龙', '万', '段', '雷', '钱', '汤', '尹', '黎', '易', '常', '武', '乔',
    '贺', '赖', '龚', '文',
]
GIVEN_CHARS = [
    '伟', '芳', '娜', '敏', '静', '丽', '强', '磊', '军', '洋', '勇', '艳',
    '杰', '娟', '涛', '明', '超', '霞', '平', '刚', '华', '建', '晓', '志',
    '俊', '浩', '然', '涵', '欣', '怡', '思', '远', '雨', '泽', '宇', '轩',
    '嘉', '梦', '琪', '雅', '婷', '慧', '东', '燕', '晨', '曦', '佳', '瑞',
    '鑫', '鹏', '飞', '雪', '梅', '兰', '虎', '龙', '文', '博', '宁', '辉',
    '颖', '妍', '薇', '倩', '楠', '璐', '玥', '悦', '昕', '彤', '瑶', '洁',
]

GRADES = ['大一', '大二', '大三', '大四', '研一', '研二']
MAJORS = [
    '计算机科学', '软件工程', '电子信息', '人工智能',
    '数据科学', '网络安全', '通信工程',
]

PROJECT_CATALOG = [
    {
        'name': '智能校园导览系统',
        'intro': '基于移动端定位与AR技术的校园导览系统，为新生和访客提供室内外一体化导航、'
                 '景点语音讲解和实时活动指引，支持多语言与无障碍模式。',
    },
    {
        'name': '基于深度学习的医学影像辅助诊断平台',
        'intro': '利用卷积神经网络对CT/MRI影像进行病灶检测与分类，辅助医生提升诊断效率与准确率，'
                 '平台提供模型训练、推理服务和医生复核工作流。',
    },
    {
        'name': '校园二手交易小程序',
        'intro': '面向在校学生的二手物品交易小程序，支持发布、求购、信用评价和线下交接预约，'
                 '保障交易安全与便捷。',
    },
    {
        'name': '农产品溯源区块链平台',
        'intro': '基于联盟链的农产品全链路溯源平台，记录种植、加工、物流、销售各环节关键数据，'
                 '消费者扫码即可查验真伪与质量信息。',
    },
    {
        'name': '智能安防监控系统',
        'intro': '结合边缘计算与视频分析的智能安防系统，实现人员闯入检测、异常行为识别和告警推送，'
                 '适用于校园和园区场景。',
    },
    {
        'name': '在线协同文档编辑器',
        'intro': '支持多人实时协同编辑的在线文档工具，基于CRDT算法实现冲突自动合并，'
                 '提供评论、版本回溯和权限管理。',
    },
    {
        'name': '基于大模型的智能客服系统',
        'intro': '集成大语言模型的智能客服系统，支持多轮对话、知识库检索增强和工单流转，'
                 '可快速接入校园各类业务咨询场景。',
    },
    {
        'name': '校园活动管理平台',
        'intro': '一站式校园活动管理平台，覆盖活动发布、报名、签到、学分认定和数据分析，'
                 '面向社团、学院和校级活动。',
    },
    {
        'name': '环境监测物联网系统',
        'intro': '基于LoRa/WiFi的环境监测物联网系统，实时采集温湿度、PM2.5、噪声等数据，'
                 '通过可视化大屏展示并支持阈值告警。',
    },
    {
        'name': '学生学业预警分析系统',
        'intro': '整合教务数据与学习行为数据，运用机器学习模型识别学业困难学生并自动预警，'
                 '辅助辅导员及时干预帮扶。',
    },
    {
        'name': '校园能耗数字孪生平台',
        'intro': '汇总楼宇水电与设备运行数据，以数字孪生方式展示能耗趋势并给出节能建议。',
    },
    {
        'name': '非遗文化数字展示平台',
        'intro': '通过三维采集、互动故事和数字档案展示地方非遗项目，支持课程实践与公众传播。',
    },
    {
        'name': '实验室预约与安全巡检系统',
        'intro': '覆盖实验室预约、准入培训、设备台账和安全巡检闭环，降低日常管理成本。',
    },
    {
        'name': '社区志愿服务协同平台',
        'intro': '连接社区需求、志愿者排班和服务评价，形成可追踪的公益服务记录。',
    },
    {
        'name': '就业岗位智能匹配助手',
        'intro': '基于学生能力画像和岗位要求进行可解释匹配，并提供简历改进与学习建议。',
    },
    {
        'name': '智慧图书馆座位与资源导航',
        'intro': '提供座位预约、馆藏定位、学习空间热度分析和无障碍导航。',
    },
    {
        'name': '城市积水风险监测终端',
        'intro': '融合低功耗传感器与气象数据，对易涝点进行实时监测和分级预警。',
    },
    {
        'name': '校园食品安全追溯平台',
        'intro': '记录供应商、检验、入库和留样信息，为师生提供透明的食品安全查询。',
    },
    {
        'name': '面向老年人的反诈训练工具',
        'intro': '以情景化案例和语音交互开展反诈训练，帮助社区老人识别常见诈骗手法。',
    },
    {
        'name': '赛事材料智能核验助手',
        'intro': '按照比赛清单自动核验申报材料完整性、版本一致性和关键字段。',
    },
    {
        'name': '开源项目贡献成长地图',
        'intro': '聚合代码、文档、评审与社区协作记录，帮助成员形成可验证的成长档案。',
    },
    {
        'name': '乡村文旅路线共创平台',
        'intro': '支持村落资源采集、路线共创、游客反馈和运营数据复盘，服务乡村文旅实践。',
    },
    {
        'name': '多模态课堂互动分析系统',
        'intro': '在获得授权的前提下分析课堂互动数据，为教学改进提供聚合指标。',
    },
    {
        'name': '科研数据合规归档工具',
        'intro': '帮助项目团队完成数据分级、版本记录、授权确认和结项归档。',
    },
]

# 小团队版的确定性项目时间线。2021 年只有三个项目；其中两个分别在
# 2022、2023 年取得国金/国银后结项。2024 年新增四个项目，与仍在持续
# 的老项目一起形成五支当前队伍。
PROJECT_DATA = (
    {
        **PROJECT_CATALOG[0],
        'code': f'{DEMO_PROJECT_PREFIX}2021-01',
        'leader_key': 'leader1',
        'start_date': date(2021, 3, 15),
        'planned_end_date': date(2022, 12, 31),
        'actual_end_date': date(2022, 12, 15),
        'status': Project.Status.CLOSED,
        'current_stage': Project.Stage.CLOSED,
        'priority': Project.Priority.NORMAL,
    },
    {
        **PROJECT_CATALOG[1],
        'code': f'{DEMO_PROJECT_PREFIX}2021-02',
        'leader_key': 'leader2',
        'start_date': date(2021, 5, 10),
        'planned_end_date': date(2023, 12, 31),
        'actual_end_date': date(2023, 12, 15),
        'status': Project.Status.CLOSED,
        'current_stage': Project.Stage.CLOSED,
        'priority': Project.Priority.NORMAL,
    },
    {
        **PROJECT_CATALOG[2],
        'code': f'{DEMO_PROJECT_PREFIX}2021-03',
        'leader_key': 'leader1',
        'start_date': date(2021, 9, 20),
        'planned_end_date': date(2027, 6, 30),
        'actual_end_date': None,
        'status': Project.Status.ACTIVE,
        'current_stage': Project.Stage.NATIONAL_COMP,
        'priority': Project.Priority.URGENT,
    },
    {
        **PROJECT_CATALOG[3],
        'code': f'{DEMO_PROJECT_PREFIX}2024-01',
        'leader_key': 'leader2',
        'start_date': date(2024, 2, 20),
        'planned_end_date': date(2027, 9, 30),
        'actual_end_date': None,
        'status': Project.Status.ACTIVE,
        'current_stage': Project.Stage.PROVINCE_COMP,
        'priority': Project.Priority.HIGH,
    },
    {
        **PROJECT_CATALOG[4],
        'code': f'{DEMO_PROJECT_PREFIX}2024-02',
        'leader_key': 'leader3',
        'start_date': date(2024, 3, 18),
        'planned_end_date': date(2027, 12, 31),
        'actual_end_date': None,
        'status': Project.Status.ACTIVE,
        'current_stage': Project.Stage.DEFENSE_PREP,
        'priority': Project.Priority.HIGH,
    },
    {
        **PROJECT_CATALOG[5],
        'code': f'{DEMO_PROJECT_PREFIX}2024-03',
        'leader_key': 'leader4',
        'start_date': date(2024, 5, 6),
        'planned_end_date': date(2028, 3, 31),
        'actual_end_date': None,
        'status': Project.Status.ACTIVE,
        'current_stage': Project.Stage.MATERIAL_SUBMIT,
        'priority': Project.Priority.NORMAL,
    },
    {
        **PROJECT_CATALOG[6],
        'code': f'{DEMO_PROJECT_PREFIX}2024-04',
        'leader_key': 'leader5',
        'start_date': date(2024, 7, 1),
        'planned_end_date': date(2028, 6, 30),
        'actual_end_date': None,
        'status': Project.Status.ACTIVE,
        'current_stage': Project.Stage.DEV_EXPERIMENT,
        'priority': Project.Priority.HIGH,
    },
)

CORE_COMPETITION_SERIES = (
    {
        'key': 'dachuang',
        'short_name': '大创',
        'name': '国家级大学生创新创业训练计划（大创）',
        'comp_type': '创新创业训练主赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '教育部',
    },
    {
        'key': 'datiao',
        'short_name': '大挑',
        'name': '挑战杯全国大学生课外学术科技作品竞赛（大挑）',
        'comp_type': '科技发明制作赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '共青团中央',
    },
    {
        'key': 'xiaotiao',
        'short_name': '小挑',
        'name': '挑战杯中国大学生创业计划竞赛（小挑）',
        'comp_type': '创业计划主赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '共青团中央',
    },
    {
        'key': 'guochuang',
        'short_name': '国创',
        'name': '中国国际大学生创新大赛（国创）',
        'comp_type': '高教主赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '教育部',
    },
)

SELECTIVE_COMPETITION_SERIES = (
    {
        'key': 'digital_china',
        'short_name': '数字中国',
        'name': '数字中国创新大赛',
        'comp_type': '数字技术创新赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '数字中国建设峰会组委会',
        'start_year': 2022,
        'entry_counts': {2022: 1, 2023: 2, 2024: 3, 2025: 2, 2026: 1},
    },
    {
        'key': 'xiamen_bank',
        'short_name': '厦门银行杯',
        'name': '厦门银行杯金融科技创新大赛',
        'comp_type': '金融科技应用赛道',
        'level': Competition.Level.NATIONAL,
        'organizer': '厦门银行杯金融科技创新大赛组委会',
        'start_year': 2023,
        'entry_counts': {2023: 1, 2024: 2, 2025: 3, 2026: 2},
    },
)

COMPETITION_DATA = CORE_COMPETITION_SERIES + SELECTIVE_COMPETITION_SERIES

# 15 个技能标签
SKILL_NAMES = [
    '前端', '后端', 'UI设计', '文档撰写', 'PPT制作', '答辩', '绘图',
    '视频制作', '数据整理', '算法', '实验', '机械设计', '财务整理',
    '比赛申报', '项目统筹',
]


def gen_chinese_name(used):
    """生成一个未使用过的中文姓名"""
    for _ in range(200):
        surname = random.choice(SURNAMES)
        length = random.randint(1, 2)
        given = ''.join(random.choice(GIVEN_CHARS) for _ in range(length))
        name = surname + given
        if name not in used:
            used.add(name)
            return name
    # 兜底：姓名 + 序号
    name = '同学' + str(len(used) + 1)
    used.add(name)
    return name


def gen_phone():
    """生成一个随机手机号（演示用）"""
    second = random.choice(['3', '5', '7', '8', '9'])
    tail = ''.join(str(random.randint(0, 9)) for _ in range(9))
    return '1' + second + tail


class Command(BaseCommand):
    """生成完整演示数据"""

    help = '生成完整的团队管理软件演示数据（账号、项目、任务、经费、知识产权等）'

    def add_arguments(self, parser):
        parser.add_argument(
            '--clean',
            action='store_true',
            dest='clean',
            default=False,
            help='先清除所有演示数据再重新生成',
        )
        parser.add_argument(
            '--force',
            action='store_true',
            dest='force',
            default=False,
            help='跳过确认提示',
        )

    # ------------------------------------------------------------------
    # 入口
    # ------------------------------------------------------------------
    def handle(self, *args, **options):
        random.seed(42)  # 固定随机种子，保证可复现

        clean = options['clean']
        force = options['force']

        # 仅以本命令专属项目编号判断是否重复运行，避免与其他演示种子互相阻塞。
        if (
            not clean
            and Project.all_objects.filter(code__startswith=DEMO_PROJECT_PREFIX).exists()
        ):
            self.stdout.write(self.style.ERROR(
                '检测到数据库中已存在团队演示数据。'
                '请使用 --clean 参数清除后重新生成。'
            ))
            return

        if not force:
            tip = '即将生成大量演示数据' + ('（会先清除现有演示数据）' if clean else '')
            confirm = input(f'{tip}，是否继续？[y/N]: ')
            if confirm.strip().lower() not in ('y', 'yes'):
                self.stdout.write(self.style.WARNING('已取消操作。'))
                return

        try:
            with transaction.atomic():
                if clean:
                    self.clean_demo_data()

                self.stdout.write(self.style.MIGRATE_HEADING('开始生成演示数据...'))

                # 容器，保存生成结果供后续步骤引用
                self.users = {}
                self.members = []
                self.leaders = []
                self.teachers = []
                self.viewing_teachers = []
                self.primary_contributors = []
                self.projects = []
                self.project_members = {}  # project -> [users]
                self.tasks_by_project = {}
                self.files_by_project = {}
                self.competitions = []
                self.competition_files_by_entry = {}

                self.create_users()
                self.create_team_organization()
                self.create_projects()
                self.create_lifecycle_history()
                self.create_competitions()
                self.create_tasks()
                self.create_finance()
                self.create_demo_files()
                self.create_import_history()
                self.create_skills()
                self.create_work_schedules()
                self.create_contributions()
                self.create_rankings()
                self.create_scheduled_reports()
                self.create_ip_applications()
                self.create_portal_publications()
                self.create_sensitive_data()
                self.create_sensitive_requests()
                self.create_announcements()
                self.create_notifications()
                self.create_operation_logs()

            self.stdout.write(self.style.SUCCESS('\n========== 演示数据生成完成 =========='))
            self.print_account_summary()
        except Exception as e:
            self.stdout.write(self.style.ERROR(f'\n生成失败，已回滚: {e}'))
            raise

    # ------------------------------------------------------------------
    # 账号
    # ------------------------------------------------------------------
    def _make_user(self, email, username, password, name, role,
                   grade='', major='', is_student=True,
                   is_staff=False, is_superuser=False, phone=''):
        user, _ = User.objects.update_or_create(
            email=email,
            defaults={
                'username': username,
                'name': name,
                'global_role': role,
                'grade': grade,
                'major': major,
                'is_student': is_student,
                'is_staff': is_staff,
                'is_superuser': is_superuser,
                'is_active': True,
                'phone': phone or gen_phone(),
                'membership_status': User.MembershipStatus.ACTIVE,
                'team_left_at': None,
                'exit_reason': '',
                'handover_to': None,
                'handover_notes': '',
            },
        )
        user.set_password(password)
        user.save()
        return user

    def create_users(self):
        self.stdout.write('-> 创建账号...')
        U = User.GlobalRole
        existing_operating_teacher = (
            User.objects.filter(
                global_role=U.TEACHER,
                is_active=True,
            )
            .exclude(membership_status=User.MembershipStatus.EXITED)
            .exclude(email__in=DEMO_ACCOUNT_EMAILS)
            .order_by('date_joined', 'id')
            .first()
        )
        # Old demo databases may still have teacher2 (or another demo account)
        # marked as a global teacher. Normalize them before assigning teacher1.
        User.objects.filter(
            email__in=DEMO_ACCOUNT_EMAILS,
            global_role=U.TEACHER,
        ).exclude(email='teacher1@demo.com').update(global_role=U.MEMBER)
        demo_teacher_role = (
            U.MEMBER
            if existing_operating_teacher
            else U.TEACHER
        )

        self.users['admin'] = self._make_user(
            'admin@demo.com', 'admin', 'admin123456', '系统管理员',
            U.SYS_ADMIN, is_student=False, is_staff=True, is_superuser=True,
            phone='13800000000',
        )
        self.users['teacher1'] = self._make_user(
            'teacher1@demo.com', 'teacher1', 'teacher123456', '张老师',
            demo_teacher_role,
            is_student=False,
            is_staff=not bool(existing_operating_teacher),
            major='计算机科学',
            phone='13800000001',
        )
        self.users['teacher2'] = self._make_user(
            'teacher2@demo.com', 'teacher2', 'teacher123456', '李老师',
            U.MEMBER, is_student=False, is_staff=False, major='软件工程',
            phone='13800000002',
        )
        self.users['teacher3'] = self._make_user(
            'teacher3@demo.com', 'teacher3', 'teacher123456', '陈老师',
            U.MEMBER, is_student=False, is_staff=False, major='工商管理',
            phone='13800000008',
        )
        self.users['teacher4'] = self._make_user(
            'teacher4@demo.com', 'teacher4', 'teacher123456', '周老师',
            U.MEMBER, is_student=False, is_staff=False, major='电子信息',
            phone='13800000009',
        )
        self.users['leader1'] = self._make_user(
            'leader1@demo.com', 'leader1', 'leader123456', '王明',
            U.MEMBER, grade='研二', major='计算机科学', phone='13800000003',
        )
        self.users['leader2'] = self._make_user(
            'leader2@demo.com', 'leader2', 'leader123456', '赵芳',
            U.MEMBER, grade='研一', major='软件工程', phone='13800000004',
        )
        self.users['leader3'] = self._make_user(
            'leader3@demo.com', 'leader3', 'leader123456', '刘强',
            U.MEMBER, grade='大三', major='电子信息', phone='13800000005',
        )
        self.users['leader4'] = self._make_user(
            'leader4@demo.com', 'leader4', 'leader123456', '陈雨桐',
            U.MEMBER, grade='研一', major='数据科学', phone='13800000007',
        )
        self.users['leader5'] = self._make_user(
            'leader5@demo.com', 'leader5', 'leader123456', '林嘉宁',
            U.MEMBER, grade='大三', major='人工智能', phone='13800000010',
        )
        self.users['contributor1'] = self._make_user(
            'contributor1@demo.com', 'contributor1', 'member123456', '钱思远',
            U.MEMBER, grade='研一', major='软件工程', phone='13800000011',
        )
        self.users['contributor2'] = self._make_user(
            'contributor2@demo.com', 'contributor2', 'member123456', '孙雅琪',
            U.MEMBER, grade='大四', major='数据科学', phone='13800000012',
        )

        self.leaders = [
            self.users['leader1'],
            self.users['leader2'],
            self.users['leader3'],
            self.users['leader4'],
            self.users['leader5'],
        ]
        self.primary_contributors = [
            self.users['contributor1'],
            self.users['contributor2'],
        ]
        self.viewing_teachers = [
            self.users['teacher2'],
            self.users['teacher3'],
            self.users['teacher4'],
        ]
        # 小团队版只保留一个全局“操作老师”。真实数据库若已有操作
        # 老师，不改动真实账号，演示老师均作为团队范围内查看老师。
        self.teachers = [
            existing_operating_teacher or self.users['teacher1']
        ]
        self.demo_teacher_is_operator = existing_operating_teacher is None

        used_names = {
            '系统管理员', '张老师', '李老师', '陈老师', '周老师',
            '王明', '赵芳', '刘强', '陈雨桐', '林嘉宁', '钱思远', '孙雅琪',
        }
        for i in range(1, 36):
            email = f'member{i}@demo.com'
            username = f'member{i}'
            name = gen_chinese_name(used_names)
            user = self._make_user(
                email, username, 'member123456', name, U.MEMBER,
                grade=random.choice(GRADES), major=random.choice(MAJORS),
            )
            self.members.append(user)
            self.users[f'member{i}'] = user

        preferred_themes = {
            'admin': 'blue',
            'teacher1': 'purple',
            'teacher2': 'green',
            'teacher3': 'blue',
            'teacher4': 'orange',
            'leader1': 'green',
            'leader2': 'blue',
            'leader3': 'purple',
            'leader4': 'orange',
            'leader5': 'green',
            'contributor1': 'blue',
            'contributor2': 'purple',
        }
        theme_cycle = ('blue', 'green', 'purple', 'orange')
        notification_categories = (
            'system',
            'task',
            'project',
            'competition',
            'finance',
            'contribution',
            'schedule',
            'approval',
            'report',
        )
        for key, user in self.users.items():
            suffix = ''.join(character for character in key if character.isdigit())
            cycle_index = int(suffix) - 1 if suffix else user.pk
            theme = preferred_themes.get(
                key, theme_cycle[cycle_index % len(theme_cycle)]
            )
            member_number = int(suffix) if key.startswith('member') else 0
            if key == 'admin':
                profile = 'admin'
                default_landing = 'dashboard'
                sidebar_order = [
                    'workspace', 'execution', 'resources',
                    'outcomes', 'administration',
                ]
                favorite_routes = [
                    '/team', '/projects', '/finance', '/reports',
                ]
                saved_filters = {
                    'projects': {'status': ['active'], 'priority': ['high', 'urgent']},
                    'approvals': {'status': ['pending']},
                    'reports': {'period': 'this_month'},
                }
                digest = 'instant'
                channels = {'in_app': True, 'email': True}
                quiet_hours = {'enabled': False, 'start': '22:00', 'end': '07:30'}
                default_scope = 'team'
            elif member_number == 35:
                profile = 'external'
                default_landing = 'tasks'
                sidebar_order = [
                    'workspace', 'execution', 'outcomes',
                    'resources', 'administration',
                ]
                favorite_routes = ['/tasks', '/files']
                saved_filters = {
                    'tasks': {'status': ['doing', 'need_help'], 'scope': 'mine'},
                    'projects': {'status': ['active'], 'scope': 'mine'},
                }
                digest = 'weekly'
                channels = {'in_app': True, 'email': False}
                quiet_hours = {'enabled': True, 'start': '21:30', 'end': '08:30'}
                default_scope = 'mine'
            elif key in {'teacher1', 'teacher2', 'teacher3', 'teacher4'}:
                profile = 'teacher'
                default_landing = 'projects'
                sidebar_order = [
                    'workspace', 'outcomes', 'execution',
                    'resources', 'administration',
                ]
                favorite_routes = [
                    '/projects', '/competitions',
                    '/intellectual-property', '/reports',
                ]
                saved_filters = {
                    'projects': {'status': ['active', 'paused']},
                    'competitions': {'status': ['preparing', 'ongoing']},
                    'approvals': {'status': ['pending']},
                }
                digest = 'daily'
                channels = {'in_app': True, 'email': True}
                quiet_hours = {'enabled': True, 'start': '22:00', 'end': '07:00'}
                default_scope = 'team'
            else:
                profile = 'member'
                default_landing = (
                    'projects' if key.startswith('leader') else 'tasks'
                )
                sidebar_order = [
                    'workspace', 'execution', 'resources',
                    'outcomes', 'administration',
                ]
                favorite_routes = ['/tasks', '/projects', '/contributions']
                saved_filters = {
                    'tasks': {'status': ['todo', 'doing'], 'scope': 'mine'},
                    'projects': {'status': ['active'], 'scope': 'mine'},
                    'contributions': {'status': ['pending', 'approved']},
                }
                digest = 'daily' if cycle_index % 2 else 'instant'
                channels = {
                    'in_app': True,
                    'email': cycle_index % 3 == 0,
                }
                quiet_hours = {
                    'enabled': True,
                    'start': '22:30',
                    'end': '07:30',
                }
                default_scope = 'mine'

            categories = {
                category: not (
                    profile == 'external'
                    and category in {'finance', 'approval', 'report'}
                )
                for category in notification_categories
            }
            UserPreference.objects.update_or_create(
                user=user,
                defaults={
                    'theme_color': theme,
                    'primary_color': UserPreference.primary_color_for_theme(theme),
                    'theme_mode': UserPreference.DEFAULT_THEME_MODE,
                    'schedule_start': UserPreference.DEFAULT_SCHEDULE_START,
                    'schedule_end': UserPreference.DEFAULT_SCHEDULE_END,
                    'dashboard_layout': {
                        'profile': profile,
                        'cards': {
                            'admin': [
                                'business', 'signals', 'priority', 'delivery',
                            ],
                            'teacher': [
                                'signals', 'business', 'delivery', 'priority',
                            ],
                            'external': [
                                'delivery', 'priority', 'signals', 'business',
                            ],
                            'member': [
                                'priority', 'delivery', 'signals', 'business',
                            ],
                        }[profile],
                    },
                    'default_landing': default_landing,
                    'sidebar_collapsed': cycle_index % 3 == 0,
                    'notification_sound': cycle_index % 4 != 0,
                    'items_per_page': (10, 20, 50)[cycle_index % 3],
                    'default_scope': default_scope,
                    'sidebar_order': sidebar_order,
                    'favorite_routes': favorite_routes,
                    'saved_filters': saved_filters,
                    'notification_preferences': {
                        'categories': categories,
                        'channels': channels,
                        'quiet_hours': quiet_hours,
                        'digest': digest,
                    },
                },
            )

        self.stdout.write(self.style.SUCCESS(
            '   账号创建完成：1 个操作老师、3 个查看老师、5 个负责人、'
            f'2 个平行主贡献者、{len(self.members)} 个普通成员'
        ))

    def create_team_organization(self):
        """创建总团队人员池；实际组队由项目和比赛参赛条目表达。"""
        self.stdout.write('-> 创建团队组织与成员关系...')
        owner = self.users['leader1']
        self.team, _ = Team.objects.update_or_create(
            code=DEMO_TEAM_CODE,
            defaults={
                'name': '数智创新实践团队',
                'description': (
                    '一个总团队人员池：1 位操作老师、3 位查看老师、'
                    '5 位项目负责人、2 位平行主贡献者与 35 位普通成员。'
                    '成员可同时参加多个“比赛届次 × 项目”参赛队。'
                ),
                'contact_email': 'teacher1@demo.com',
                'join_message': (
                    '欢迎愿意持续投入真实项目、遵守协作规范并主动复盘的成员加入。'
                ),
                'is_active': True,
                'owner': owner,
            },
        )
        team_created_at = timezone.make_aware(
            datetime.combine(date(2021, 1, 10), time(9, 0))
        )
        Team.objects.filter(pk=self.team.pk).update(created_at=team_created_at)

        memberships = {}
        ordered_keys = [
            'teacher1', 'teacher2', 'teacher3', 'teacher4',
            'leader1', 'leader2', 'leader3', 'leader4', 'leader5',
            'contributor1', 'contributor2',
            *(f'member{number}' for number in range(1, 36)),
        ]
        ordered_users = [(key, self.users[key]) for key in ordered_keys]
        for index, (key, user) in enumerate(ordered_users):
            member_number = (
                int(key.removeprefix('member'))
                if key.startswith('member')
                else 0
            )
            if key == 'leader1':
                role = TeamMember.Role.OWNER
            elif key in {'teacher1', 'teacher2', 'teacher3', 'teacher4'}:
                role = TeamMember.Role.TEACHER
            elif key in {
                'leader2', 'leader3', 'leader4', 'leader5',
                'contributor1', 'contributor2',
            }:
                role = TeamMember.Role.CO_LEAD
            elif member_number == 35:
                role = TeamMember.Role.EXTERNAL
            else:
                role = TeamMember.Role.MEMBER

            if 30 <= member_number <= 32:
                status = TeamMember.Status.ON_LEAVE
            elif 33 <= member_number <= 34:
                status = TeamMember.Status.EXITED
            else:
                status = TeamMember.Status.ACTIVE

            joined_at = team_created_at + timedelta(days=index * 11)
            membership = TeamMember.objects.create(
                team=self.team,
                user=user,
                role=role,
                status=status,
                left_at=(
                    joined_at + timedelta(days=760)
                    if status == TeamMember.Status.EXITED
                    else None
                ),
                exit_reason=(
                    f'{DEMO_MARKER}完成阶段性学习与项目交接'
                    if status == TeamMember.Status.EXITED
                    else ''
                ),
                handover_notes=(
                    f'{DEMO_MARKER}任务、文件与项目进展已完成交接'
                    if status == TeamMember.Status.EXITED
                    else ''
                ),
            )
            TeamMember.objects.filter(pk=membership.pk).update(joined_at=joined_at)
            membership.joined_at = joined_at
            memberships[key] = membership

            event = TeamMembershipEvent.objects.create(
                membership=membership,
                event_type='joined',
                to_role=role,
                to_status=TeamMember.Status.ACTIVE,
                reason=f'{DEMO_MARKER}加入数智创新实践团队',
                operator=owner,
            )
            TeamMembershipEvent.objects.filter(pk=event.pk).update(
                created_at=joined_at
            )

        handover_target = memberships['leader2']
        for member_number in range(30, 35):
            key = f'member{member_number}'
            membership = memberships[key]
            change_at = membership.joined_at + timedelta(days=700)
            event_type = (
                'status_changed'
                if membership.status == TeamMember.Status.ON_LEAVE
                else 'exited'
            )
            event = TeamMembershipEvent.objects.create(
                membership=membership,
                event_type=event_type,
                from_status=TeamMember.Status.ACTIVE,
                to_status=membership.status,
                reason=(
                    f'{DEMO_MARKER}学业安排暂离'
                    if membership.status == TeamMember.Status.ON_LEAVE
                    else f'{DEMO_MARKER}完成阶段性学习并离队'
                ),
                operator=owner,
            )
            TeamMembershipEvent.objects.filter(pk=event.pk).update(
                created_at=change_at
            )

            if membership.status == TeamMember.Status.EXITED:
                membership.handover_to = handover_target
                membership.save(update_fields=['handover_to'])
                handover_event = TeamMembershipEvent.objects.create(
                    membership=membership,
                    event_type='handover',
                    from_status=TeamMember.Status.ACTIVE,
                    to_status=TeamMember.Status.EXITED,
                    reason=f'{DEMO_MARKER}离队交接',
                    handover_to=handover_target,
                    handover_notes=membership.handover_notes,
                    operator=owner,
                )
                TeamMembershipEvent.objects.filter(
                    pk=handover_event.pk
                ).update(created_at=change_at + timedelta(minutes=5))

        self.stdout.write(self.style.SUCCESS(
            '   总团队创建完成：成员 46 人，成员关系事件 53 条；'
            '不创建固定部门，组队关系由项目和比赛条目承载'
        ))

    # ------------------------------------------------------------------
    # 项目
    # ------------------------------------------------------------------
    def create_projects(self):
        self.stdout.write('-> 创建项目...')
        now = timezone.now()

        for idx, pdata in enumerate(PROJECT_DATA):
            leader = self.users[pdata['leader_key']]
            start_date = pdata['start_date']
            actual_end = pdata['actual_end_date']
            if actual_end:
                last_update = timezone.make_aware(
                    datetime.combine(actual_end, time(17, 30))
                )
            else:
                last_update = now - timedelta(days=(idx * 2) % 11)

            project = Project.objects.create(
                name=pdata['name'],
                code=pdata['code'],
                leader=leader,
                current_stage=pdata['current_stage'],
                start_date=start_date,
                planned_end_date=pdata['planned_end_date'],
                actual_end_date=actual_end,
                status=pdata['status'],
                priority=pdata['priority'],
                intro=pdata['intro'],
                last_leader_update=last_update,
                visibility=Project.Visibility.PROJECT,
            )
            created_at = timezone.make_aware(
                datetime.combine(start_date, time(9, 0))
            )
            update_fields = {'created_at': created_at, 'updated_at': created_at}
            if actual_end:
                update_fields['archived_at'] = timezone.make_aware(
                    datetime.combine(actual_end, time(18, 0))
                )
            Project.all_objects.filter(pk=project.pk).update(**update_fields)
            self.projects.append(project)
            # 根团队是人员池；具体组队由项目成员和比赛参赛名单表达。
            project.teams.set([self.team])

            # 历史项目 8 人、当前项目 10 人。两位平行主贡献者跨项目参与，
            # 普通成员按滑动窗口分配，因此一人可以在多个队伍中协作。
            contributor_count = 1 if actual_end else 2
            chosen_contributors = [
                self.primary_contributors[
                    (idx + offset) % len(self.primary_contributors)
                ]
                for offset in range(contributor_count)
            ]
            regular_count = 6 if actual_end else 7
            chosen_regular_members = [
                self.members[(idx * 5 + offset) % 29]
                for offset in range(regular_count)
            ]
            chosen = list(dict.fromkeys(
                chosen_contributors + chosen_regular_members
            ))
            member_users = [leader, *chosen]

            self.project_members[project.id] = member_users

            # 负责人作为 leader 角色
            ProjectMember.objects.create(
                project=project, user=leader,
                role_in_project=ProjectMember.RoleInProject.LEADER,
            )
            for u in chosen:
                role = (
                    ProjectMember.RoleInProject.CORE
                    if u in self.primary_contributors
                    else ProjectMember.RoleInProject.PARTICIPANT
                )
                ProjectMember.objects.create(
                    project=project, user=u, role_in_project=role,
                )

        self.stdout.write(self.style.SUCCESS(
            f'   项目创建完成：{len(self.projects)} 个'
        ))

    def create_lifecycle_history(self):
        """生成团队成员、项目成员和项目阶段的真实时间跨度记录。"""
        self.stdout.write('-> 创建成员与项目生命周期...')
        lifecycle_users = self.leaders + self.primary_contributors + self.members

        for index, user in enumerate(lifecycle_users):
            joined_year = min(2021 + index // 8, 2026)
            joined_date = date(joined_year, 1 + index % 10, 3 + index % 20)
            status = User.MembershipStatus.ACTIVE
            team_left_at = None
            exit_reason = ''
            handover_to = None
            handover_notes = ''

            member_number = (
                int(user.username.removeprefix('member'))
                if user.username.startswith('member')
                else 0
            )
            if 30 <= member_number <= 32:
                status = User.MembershipStatus.ON_LEAVE
            elif 33 <= member_number <= 34:
                status = User.MembershipStatus.EXITED
                team_left_at = timezone.now() - timedelta(
                    days=20 + member_number
                )
                exit_reason = '学业阶段变化，完成资料与任务交接后离队。'
                handover_to = self.leaders[member_number % len(self.leaders)]
                handover_notes = '代码仓库、项目材料与未结任务均已完成清单式交接。'
            elif member_number == 35:
                status = User.MembershipStatus.EXTERNAL

            User.objects.filter(pk=user.pk).update(
                team_joined_at=joined_date,
                membership_status=status,
                is_active=status != User.MembershipStatus.EXITED,
                team_left_at=team_left_at,
                exit_reason=exit_reason,
                handover_to=handover_to,
                handover_notes=handover_notes,
            )
            user.team_joined_at = joined_date
            user.membership_status = status
            user.is_active = status != User.MembershipStatus.EXITED

            joined_event = UserLifecycleEvent.objects.create(
                user=user,
                event_type=UserLifecycleEvent.EventType.CREATED,
                to_status=User.MembershipStatus.ACTIVE,
                reason=f'{DEMO_MARKER}经面试与试做任务后加入团队。',
                operator=self.teachers[index % len(self.teachers)],
            )
            joined_at = timezone.make_aware(
                datetime.combine(joined_date, time(10, 0))
            )
            UserLifecycleEvent.objects.filter(pk=joined_event.pk).update(
                created_at=joined_at
            )

            if status != User.MembershipStatus.ACTIVE:
                change_event = UserLifecycleEvent.objects.create(
                    user=user,
                    event_type=UserLifecycleEvent.EventType.STATUS_CHANGED,
                    from_status=User.MembershipStatus.ACTIVE,
                    to_status=status,
                    reason=f'{DEMO_MARKER}{exit_reason or "根据学习安排调整团队参与状态。"}',
                    handover_to=handover_to,
                    handover_notes=handover_notes,
                    operator=self.teachers[index % len(self.teachers)],
                )
                event_time = team_left_at or (timezone.now() - timedelta(days=7 + index))
                UserLifecycleEvent.objects.filter(pk=change_event.pk).update(
                    created_at=event_time
                )

        stage_flow = [
            Project.Stage.CONCEIVING,
            Project.Stage.APPROVED,
            Project.Stage.MATERIAL_PREP,
            Project.Stage.DEV_EXPERIMENT,
            Project.Stage.MATERIAL_SUBMIT,
            Project.Stage.DEFENSE_PREP,
            Project.Stage.SCHOOL_COMP,
            Project.Stage.PROVINCE_COMP,
            Project.Stage.NATIONAL_COMP,
            Project.Stage.AWARDED,
            Project.Stage.CLOSED,
        ]
        membership_event_count = 0
        stage_log_count = 0

        for project_index, project in enumerate(self.projects):
            project_start = timezone.make_aware(
                datetime.combine(project.start_date, time(9, 0))
            )
            if project.status == Project.Status.PAUSED:
                stages = stage_flow[:4] + [Project.Stage.PAUSED]
            else:
                stages = [
                    stage for stage in stage_flow
                    if stage <= project.current_stage
                ]
            previous_stage = None
            for stage_index, stage in enumerate(stages):
                log = ProjectStageLog.objects.create(
                    project=project,
                    from_stage=previous_stage,
                    to_stage=stage,
                    operator=project.leader,
                    note=f'{DEMO_MARKER}{project.code} 完成“{Project.Stage(stage).label}”节点。',
                )
                ProjectStageLog.objects.filter(pk=log.pk).update(
                    created_at=project_start + timedelta(days=stage_index * 35)
                )
                previous_stage = stage
                stage_log_count += 1

            memberships = list(
                ProjectMember.objects.filter(project=project).order_by('id')
            )
            for member_index, membership in enumerate(memberships):
                joined_at = project_start + timedelta(days=member_index * 3)
                ProjectMember.objects.filter(pk=membership.pk).update(
                    joined_at=joined_at
                )
                joined = ProjectMembershipEvent.objects.create(
                    membership=membership,
                    event_type=ProjectMembershipEvent.EventType.JOINED,
                    to_role=membership.role_in_project,
                    to_status=ProjectMember.Status.ACTIVE,
                    reason=f'{DEMO_MARKER}按项目阶段需要加入团队。',
                    operator=project.leader,
                )
                ProjectMembershipEvent.objects.filter(pk=joined.pk).update(
                    created_at=joined_at
                )
                membership_event_count += 1

            candidates = [m for m in memberships if m.user_id != project.leader_id]
            if not candidates:
                continue
            target = candidates[-1]
            handover = candidates[0] if candidates[0].pk != target.pk else memberships[0]
            if project_index % 3 == 0:
                exited_at = min(
                    timezone.now() - timedelta(days=10 + project_index),
                    project_start + timedelta(days=180),
                )
                ProjectMember.objects.filter(pk=target.pk).update(
                    status=ProjectMember.Status.EXITED,
                    exited_at=exited_at,
                    exit_reason='课程与实习安排变化，退出当前项目。',
                    handover_to=handover,
                    handover_notes='已交接材料目录、任务进度和后续联系人。',
                )
                event = ProjectMembershipEvent.objects.create(
                    membership=target,
                    event_type=ProjectMembershipEvent.EventType.EXITED,
                    from_status=ProjectMember.Status.ACTIVE,
                    to_status=ProjectMember.Status.EXITED,
                    reason=f'{DEMO_MARKER}课程与实习安排变化，退出当前项目。',
                    handover_to=handover,
                    handover_notes='已交接材料目录、任务进度和后续联系人。',
                    operator=project.leader,
                )
                ProjectMembershipEvent.objects.filter(pk=event.pk).update(
                    created_at=exited_at
                )
                membership_event_count += 1
            elif project_index % 3 == 1:
                ProjectMember.objects.filter(pk=target.pk).update(
                    status=ProjectMember.Status.ON_LEAVE,
                    exit_reason='考试周期间暂离项目。',
                )
                event = ProjectMembershipEvent.objects.create(
                    membership=target,
                    event_type=ProjectMembershipEvent.EventType.STATUS_CHANGED,
                    from_status=ProjectMember.Status.ACTIVE,
                    to_status=ProjectMember.Status.ON_LEAVE,
                    reason=f'{DEMO_MARKER}考试周期间暂离项目，结束后恢复参与。',
                    operator=project.leader,
                )
                ProjectMembershipEvent.objects.filter(pk=event.pk).update(
                    created_at=timezone.now() - timedelta(days=5 + project_index)
                )
                membership_event_count += 1

        self.stdout.write(self.style.SUCCESS(
            f'   生命周期创建完成：成员事件 {UserLifecycleEvent.objects.filter(reason__startswith=DEMO_MARKER).count()} 条，'
            f'项目成员事件 {membership_event_count} 条，阶段日志 {stage_log_count} 条'
        ))

    # ------------------------------------------------------------------
    # 比赛
    # ------------------------------------------------------------------
    def create_competitions(self):
        """创建 2021—2026 逐届赛事、项目参赛条目和真实参赛名单。"""
        self.stdout.write('-> 创建逐届比赛、参赛队与名单...')
        current_year = timezone.localdate().year
        project_index_by_id = {
            project.id: index for index, project in enumerate(self.projects)
        }
        event_count = 0
        participant_count = 0
        award_count = 0

        def active_projects(year):
            return [
                project
                for project in self.projects
                if project.start_date.year <= year
                and (
                    project.actual_end_date is None
                    or project.actual_end_date.year >= year
                )
            ]

        def create_event_entries(series, year, projects):
            nonlocal event_count, participant_count, award_count
            event = CompetitionEvent.objects.create(
                organization=self.team,
                name=series['name'],
                edition=str(year),
                organizer=series['organizer'],
            )
            event_created_at = timezone.make_aware(
                datetime.combine(date(year, 1, 15), time(9, 0))
            )
            CompetitionEvent.objects.filter(pk=event.pk).update(
                created_at=event_created_at,
                updated_at=event_created_at,
            )
            event_count += 1

            for project in projects:
                project_index = project_index_by_id[project.id]
                is_current_edition = year == current_year
                status = (
                    Competition.Status.ONGOING
                    if is_current_edition
                    else Competition.Status.COMPLETED
                )
                result_date = date(year, 11, 15)
                award_level = ''
                if (
                    series['key'] == 'guochuang'
                    and year == 2022
                    and project.code == f'{DEMO_PROJECT_PREFIX}2021-01'
                ):
                    award_level = '国赛金奖'
                elif (
                    series['key'] == 'datiao'
                    and year == 2023
                    and project.code == f'{DEMO_PROJECT_PREFIX}2021-02'
                ):
                    award_level = '国赛银奖'
                elif (
                    not is_current_edition
                    and (year + project_index + len(series['key'])) % 5 == 0
                ):
                    award_level = '国赛铜奖'

                competition = Competition.objects.create(
                    project=project,
                    event=event,
                    entry_name=(
                        f'{project.name}·{series["short_name"]}{year}参赛队'
                    ),
                    name=series['name'],
                    comp_type=series['comp_type'],
                    level=series['level'],
                    organizer=series['organizer'],
                    register_date=date(year, 3, 15),
                    material_deadline=date(year, 5, 20),
                    review_date=date(year, 7, 10),
                    defense_date=date(year, 9, 20),
                    school_date=date(year, 4, 15),
                    province_date=date(year, 7, 25),
                    national_date=date(year, 10, 20),
                    result_date=(
                        None if is_current_edition else result_date
                    ),
                    status=status,
                    is_promoted=not is_current_edition,
                    is_awarded=bool(award_level),
                    award_level=award_level,
                    review_summary=(
                        '完成本届材料、网评与答辩复盘，相关版本均已归档。'
                        if not is_current_edition
                        else ''
                    ),
                    current_stage=(
                        '全国赛备赛'
                        if is_current_edition
                        else '本届赛事已归档'
                    ),
                )
                competition_created_at = timezone.make_aware(
                    datetime.combine(date(year, 3, 15), time(10, 0))
                )
                Competition.objects.filter(pk=competition.pk).update(
                    created_at=competition_created_at,
                    updated_at=competition_created_at,
                )
                self.competitions.append(competition)

                project_members = [
                    user
                    for user in self.project_members[project.id]
                    if user.id != project.leader_id
                ]
                selected_members = project_members[:4]
                advisor = self.users[
                    f'teacher{1 + (len(self.competitions) - 1) % 4}'
                ]
                participant_specs = [
                    (
                        project.leader,
                        CompetitionParticipant.Role.LEADER,
                        '统筹本届比赛、确认分工并负责最终提交。',
                    ),
                    *[
                        (
                            user,
                            CompetitionParticipant.Role.MEMBER,
                            (
                                '负责本届计划书、数据、系统或路演材料中的'
                                f'第 {member_index + 1} 项工作。'
                            ),
                        )
                        for member_index, user in enumerate(selected_members)
                    ],
                    (
                        advisor,
                        CompetitionParticipant.Role.ADVISOR,
                        '查看关键节点并提供方向性指导。',
                    ),
                ]
                participant_records = []
                for participant_index, (user, role, responsibility) in enumerate(
                    participant_specs
                ):
                    participant = CompetitionParticipant.objects.create(
                        competition=competition,
                        user=user,
                        role=role,
                        participation_status=(
                            CompetitionParticipant.ParticipationStatus.CONFIRMED
                        ),
                        responsibility=responsibility,
                    )
                    joined_at = competition_created_at + timedelta(
                        hours=participant_index
                    )
                    CompetitionParticipant.objects.filter(
                        pk=participant.pk
                    ).update(joined_at=joined_at, updated_at=joined_at)
                    participant_records.append(participant)
                    participant_count += 1

                if award_level:
                    award = CompetitionAward.objects.create(
                        competition=competition,
                        award_name=f'{series["short_name"]}{year}年度获奖',
                        award_level=award_level,
                        award_date=result_date,
                        notes=(
                            '获奖人取自该参赛条目的已确认名单；'
                            '国金/国银项目在同赛事同赛道不再参加后续届次。'
                        ),
                    )
                    award.recipients.set([
                        participant.user
                        for participant in participant_records[:3]
                    ])
                    award_count += 1

        for year in range(2021, current_year + 1):
            year_projects = active_projects(year)
            for series in CORE_COMPETITION_SERIES:
                create_event_entries(series, year, year_projects)

        for series_index, series in enumerate(SELECTIVE_COMPETITION_SERIES):
            for year in range(series['start_year'], current_year + 1):
                year_projects = active_projects(year)
                requested_count = series['entry_counts'].get(year, 1)
                count = min(max(1, requested_count), 3, len(year_projects))
                offset = (year + series_index) % len(year_projects)
                selected = [
                    year_projects[(offset + index) % len(year_projects)]
                    for index in range(count)
                ]
                create_event_entries(series, year, selected)

        self.stdout.write(self.style.SUCCESS(
            f'   比赛创建完成：届次 {event_count} 个、参赛条目 '
            f'{len(self.competitions)} 个、参赛人次 {participant_count}、'
            f'获奖记录 {award_count} 条'
        ))

    # ------------------------------------------------------------------
    # 任务
    # ------------------------------------------------------------------
    def create_tasks(self):
        self.stdout.write('-> 创建任务...')
        now = timezone.now()
        task_titles = [
            '需求调研与整理', '原型设计', '数据库表结构设计', '接口文档编写',
            '前端页面开发', '后端接口开发', '核心算法实现', '单元测试编写',
            '系统集成联调', '用户手册撰写', '答辩PPT制作', '比赛报名材料准备',
            '软件说明书撰写', '源代码文档整理', '性能优化', '部署上线',
            'Bug 修复与回归', '数据采集与清洗', '模型训练与调优', 'UI 视觉规范制定',
        ]
        active_statuses = [
            Task.Status.TODO,
            Task.Status.DOING,
            Task.Status.PENDING_REVIEW,
            Task.Status.DONE,
            Task.Status.OVERDUE,
            Task.Status.NEED_HELP,
        ]
        total = 0
        for pidx, project in enumerate(self.projects):
            members = self.project_members[project.id]
            leader = project.leader
            project_tasks = []
            project_started_at = timezone.make_aware(
                datetime.combine(project.start_date, time(9, 0))
            )
            for task_index in range(5):
                if project.status == Project.Status.CLOSED:
                    status = (
                        Task.Status.CANCELLED
                        if task_index == 4 and pidx % 2
                        else Task.Status.DONE
                    )
                elif project.status == Project.Status.PAUSED:
                    status = (
                        Task.Status.PAUSED if task_index >= 3 else Task.Status.DONE
                    )
                else:
                    status = active_statuses[
                        (pidx + task_index) % len(active_statuses)
                    ]

                assignee = members[(pidx + task_index) % len(members)]
                title = task_titles[
                    (pidx * 5 + task_index) % len(task_titles)
                ]
                title = f'{title}（{project.code}）'
                deadline = None
                completed_at = None
                overdue_reminded = False
                task_created_at = min(
                    now - timedelta(days=1),
                    project_started_at + timedelta(days=20 + task_index * 28),
                )
                started_at = task_created_at + timedelta(days=1)

                if status in (
                    Task.Status.TODO,
                    Task.Status.DOING,
                    Task.Status.NEED_HELP,
                ):
                    deadline = now + timedelta(days=random.randint(2, 14))
                elif status == Task.Status.PENDING_REVIEW:
                    deadline = now - timedelta(days=1)
                elif status == Task.Status.DONE:
                    completed_at = min(
                        now - timedelta(days=1),
                        task_created_at + timedelta(days=12),
                    )
                    deadline = completed_at + timedelta(days=2)
                elif status == Task.Status.OVERDUE:
                    deadline = now - timedelta(days=3 + task_index)
                    overdue_reminded = True
                elif status == Task.Status.PAUSED:
                    deadline = now + timedelta(days=30)
                elif status == Task.Status.CANCELLED:
                    deadline = task_created_at + timedelta(days=10)

                reviewer = (
                    leader
                    if status in (Task.Status.PENDING_REVIEW, Task.Status.DONE)
                    else None
                )
                completion_note = ''
                if status == Task.Status.PENDING_REVIEW:
                    completion_note = (
                        '本轮交付物与说明已上传，关键检查项已自测通过，'
                        '现提交负责人复核。'
                    )
                elif status == Task.Status.DONE:
                    completion_note = (
                        '交付物已完成并通过负责人验收，相关文件已归入项目档案。'
                    )
                task = Task.objects.create(
                    project=project,
                    title=title,
                    assignee=assignee,
                    creator=leader,
                    status=status,
                    priority=(
                        Task.Priority.HIGH
                        if task_index in (1, 4)
                        else Task.Priority.MEDIUM
                    ),
                    start_date=started_at,
                    deadline=deadline,
                    completed_at=completed_at,
                    overdue_reminded=overdue_reminded,
                    reviewer=reviewer,
                    description=f'本项目「{project.name}」的{title}任务，由{assignee.name}负责。',
                    delay_reason=(
                        '等待外部数据确认，已同步负责人。'
                        if status in (Task.Status.OVERDUE, Task.Status.NEED_HELP)
                        else ''
                    ),
                    completion_note=completion_note,
                )
                Task.all_objects.filter(pk=task.pk).update(
                    created_at=task_created_at,
                    updated_at=completed_at or task_created_at,
                )
                collaborator = members[
                    (pidx + task_index + 1) % len(members)
                ]
                if collaborator.pk != assignee.pk:
                    task.collaborators.add(collaborator)
                project_tasks.append(task)
                total += 1
            self.tasks_by_project[project.id] = project_tasks

        self.stdout.write(self.style.SUCCESS(f'   任务创建完成：{total} 个'))

    # ------------------------------------------------------------------
    # 经费
    # ------------------------------------------------------------------
    @staticmethod
    def _build_receipt_png(project, expense, receipt_index):
        """生成可被 Pillow 与 OCR 校验的简洁英文票据图片。"""
        image = PILImage.new('RGB', (960, 420), color='#ffffff')
        draw = ImageDraw.Draw(image)
        font_candidates = (
            Path('C:/Windows/Fonts/arial.ttf'),
            Path('/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'),
        )
        font_path = next((path for path in font_candidates if path.exists()), None)
        font = (
            ImageFont.truetype(str(font_path), 38)
            if font_path
            else ImageFont.load_default()
        )
        small_font = (
            ImageFont.truetype(str(font_path), 28)
            if font_path
            else ImageFont.load_default()
        )
        draw.rounded_rectangle(
            (28, 28, 932, 392),
            radius=18,
            outline='#176b73',
            width=4,
            fill='#f8fbfb',
        )
        draw.text((72, 72), 'DEMO SUPPLY STORE', fill='#17353a', font=font)
        draw.text(
            (72, 165),
            f'DATE: {expense.expense_date.isoformat()}    '
            f'NO: {project.code[-10:]}-{expense.pk}-{receipt_index}',
            fill='#334e52',
            font=small_font,
        )
        draw.text(
            (72, 260),
            f'TOTAL: {expense.amount:.2f}',
            fill='#176b73',
            font=font,
        )
        buffer = BytesIO()
        image.save(buffer, format='PNG', optimize=True)
        return buffer.getvalue()

    def create_finance(self):
        self.stdout.write('-> 创建经费数据...')
        now = timezone.now()
        budget_count = 0
        income_count = 0
        expense_count = 0
        receipt_count = 0

        expense_titles = {
            'material': '实验材料采购',
            'equipment': '开发设备购置',
            'printing': '项目材料打印',
            'travel': '调研差旅费',
            'software': '开发软件授权',
            'competition_fee': '比赛报名费',
            'promotion': '项目宣传推广',
            'labor': '劳务支出',
            'other': '其他杂项支出',
        }

        for project in self.projects:
            members = self.project_members[project.id]
            project_entries = list(
                Competition.objects.filter(project=project)
                .select_related('event')
                .order_by('-event__edition', 'event__name', 'id')
            )
            awarded_entry = next(
                (entry for entry in project_entries if entry.is_awarded),
                None,
            )
            if awarded_entry is None or len(project_entries) < 2:
                raise CommandError(
                    f'项目 {project.code} 缺少获奖或多比赛参赛条目，'
                    '无法生成可追溯经费演示数据。'
                )

            current_edition = str(timezone.localdate().year)
            current_entries = [
                entry
                for entry in project_entries
                if entry.event.edition == current_edition
            ]
            trace_entries = (
                current_entries
                if len(current_entries) >= 2
                else project_entries
            )[:2]
            competition_expense_specs = (
                {
                    'category': FinanceExpense.Category.TRAVEL,
                    'title': '比赛现场往返交通费',
                    'competition_entry': trace_entries[0],
                    'purpose': (
                        f'参加 {trace_entries[0].event.name}'
                        f'（{trace_entries[0].event.edition}）现场比赛的'
                        '往返交通与市内接驳。'
                    ),
                },
                {
                    'category': FinanceExpense.Category.COMPETITION_FEE,
                    'title': '比赛报名费',
                    'competition_entry': trace_entries[1],
                    'purpose': (
                        f'缴纳 {trace_entries[1].event.name}'
                        f'（{trace_entries[1].event.edition}）参赛报名费用。'
                    ),
                },
            )

            bonus = Decimal(random.randint(5000, 50000))
            other_income = Decimal(random.randint(1000, 10000))
            FinanceIncome.objects.create(
                project=project,
                competition_entry=awarded_entry,
                title='比赛奖金入账',
                amount=bonus,
                income_type=FinanceIncome.IncomeType.BONUS,
                income_date=(now - timedelta(days=random.randint(30, 90))).date(),
                source='赛事主办方',
                reference_number=f'DEMO-BONUS-{project.id:04d}',
                recorded_by=project.leader,
            )
            income_count += 1
            FinanceIncome.objects.create(
                project=project,
                title='项目配套经费',
                amount=other_income,
                income_type=FinanceIncome.IncomeType.GRANT,
                income_date=(now - timedelta(days=random.randint(10, 45))).date(),
                source='团队项目经费',
                reference_number=f'DEMO-GRANT-{project.id:04d}',
                recorded_by=project.leader,
            )
            income_count += 1

            # 生成 2~3 条支出，覆盖草稿、待审核、已审核、已付款状态。
            expense_num = random.randint(2, 3)
            workflow_statuses = [
                FinanceExpense.ReimbursementStatus.PAID,
                FinanceExpense.ReimbursementStatus.PENDING,
                FinanceExpense.ReimbursementStatus.APPROVED,
            ]
            for expense_index in range(expense_num):
                expense_spec = (
                    competition_expense_specs[expense_index]
                    if expense_index < len(competition_expense_specs)
                    else None
                )
                category = (
                    expense_spec['category']
                    if expense_spec
                    else random.choice(list(expense_titles.keys()))
                )
                title = (
                    expense_spec['title']
                    if expense_spec
                    else expense_titles[category]
                )
                competition_entry = (
                    expense_spec['competition_entry']
                    if expense_spec
                    else None
                )
                purpose = (
                    expense_spec['purpose']
                    if expense_spec
                    else f'{project.name} - {title}，用于项目推进所需。'
                )
                amount = Decimal(random.randint(100, 5000))
                spender = random.choice(members)
                expense_date = (now - timedelta(days=random.randint(1, 60))).date()
                reimbursement_status = workflow_statuses[expense_index % len(workflow_statuses)]
                expense = FinanceExpense.objects.create(
                    project=project,
                    competition_entry=competition_entry,
                    title=title,
                    amount=amount,
                    spender=spender,
                    expense_date=expense_date,
                    category=category,
                    purpose=purpose,
                    reimbursement_status=reimbursement_status,
                    applied_by=spender,
                    applied_at=now - timedelta(days=max(1, 8 - expense_index)),
                    reviewer=project.leader if reimbursement_status != FinanceExpense.ReimbursementStatus.PENDING else None,
                    reviewed_at=now - timedelta(days=max(1, 6 - expense_index)) if reimbursement_status != FinanceExpense.ReimbursementStatus.PENDING else None,
                    review_opinion='票据与用途核验通过。' if reimbursement_status != FinanceExpense.ReimbursementStatus.PENDING else '',
                    paid_by=project.leader if reimbursement_status == FinanceExpense.ReimbursementStatus.PAID else None,
                    paid_at=now - timedelta(days=2) if reimbursement_status == FinanceExpense.ReimbursementStatus.PAID else None,
                    payment_method='银行转账' if reimbursement_status == FinanceExpense.ReimbursementStatus.PAID else '',
                    payment_reference=f'DEMO-PAY-{project.id:04d}-{expense_index + 1}' if reimbursement_status == FinanceExpense.ReimbursementStatus.PAID else '',
                )
                expense_count += 1

                # 为每条支出创建 1~2 张模拟票据
                for _r in range(random.randint(1, 2)):
                    receipt = FinanceReceipt(
                        expense=expense,
                        uploaded_by=spender,
                    )
                    receipt_data = self._build_receipt_png(
                        project,
                        expense,
                        _r + 1,
                    )
                    receipt.file.save(
                        f'receipt_{project.id}_{expense.id}_{_r+1}.png',
                        ContentFile(receipt_data),
                        save=True,
                    )
                    receipt_count += 1

            budget = FinanceBudget.objects.filter(project=project).order_by('-updated_at').first()
            budget.period = now.strftime('%Y-%m')
            budget.save(update_fields=['period', 'updated_at'])
            budget_count += 1

        self.stdout.write(self.style.SUCCESS(
            f'   经费创建完成：收入 {income_count} 条，预算 {budget_count} 条，'
            f'支出 {expense_count} 条，票据 {receipt_count} 条'
        ))

    def backend_dir(self):
        return Path(__file__).resolve().parents[4]

    def assets_dir(self):
        return self.backend_dir() / 'seed_assets' / 'competition_demo_files'

    def _project_document_lines(self, project):
        return [
            f'项目名称：{project.name}',
            f'项目编号：{project.code}',
            f'项目负责人：{project.leader.name}',
            f'项目摘要：{project.intro[:90]}',
            '说明：本文件由团队演示数据命令按当前项目动态生成。',
        ]

    @staticmethod
    def _pdf_font_name():
        """注册能嵌入中文并保留文本映射的 PDF 字体。"""
        font_name = 'SeedDemoCJK'
        if font_name in pdfmetrics.getRegisteredFontNames():
            return font_name
        candidates = (
            Path('C:/Windows/Fonts/simhei.ttf'),
            Path('C:/Windows/Fonts/simsunb.ttf'),
            Path('/usr/share/fonts/truetype/wqy/wqy-zenhei.ttc'),
            Path('/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc'),
            Path('/usr/share/fonts/truetype/arphic/uming.ttc'),
        )
        for candidate in candidates:
            if not candidate.exists():
                continue
            try:
                pdfmetrics.registerFont(TTFont(font_name, str(candidate)))
                return font_name
            except Exception:
                continue

        # 最后使用 ReportLab 内置中文字体，保证最小环境仍能生成合法 PDF。
        fallback_name = 'STSong-Light'
        if fallback_name not in pdfmetrics.getRegisteredFontNames():
            pdfmetrics.registerFont(UnicodeCIDFont(fallback_name))
        return fallback_name

    def _build_project_pdf(self, project):
        buffer = BytesIO()
        document = canvas.Canvas(buffer, pagesize=A4)
        document.setTitle(f'{project.code} {project.name}')
        document.setAuthor(project.leader.name)
        font_name = self._pdf_font_name()
        document.setFont(font_name, 20)
        document.drawString(54, 790, '团队项目概览')
        document.setFont(font_name, 12)
        y = 744
        for line in self._project_document_lines(project):
            for segment_start in range(0, len(line), 38):
                document.drawString(
                    54,
                    y,
                    line[segment_start:segment_start + 38],
                )
                y -= 25
            y -= 8
        document.save()
        return buffer.getvalue()

    def _build_project_docx(self, project):
        document = Document()
        document.add_heading('团队项目工作计划', level=1)
        for line in self._project_document_lines(project):
            document.add_paragraph(line)
        document.add_paragraph('下一步：完成当前阶段任务并更新项目进展。')
        buffer = BytesIO()
        document.save(buffer)
        return buffer.getvalue()

    def _build_project_xlsx(self, project):
        workbook = Workbook()
        worksheet = workbook.active
        worksheet.title = '项目简表'
        worksheet.append(['字段', '内容'])
        worksheet.append(['项目名称', project.name])
        worksheet.append(['项目编号', project.code])
        worksheet.append(['项目负责人', project.leader.name])
        worksheet.append(['项目摘要', project.intro[:90]])
        worksheet.append(['数据说明', '本工作簿由团队演示数据命令动态生成'])
        worksheet.column_dimensions['A'].width = 18
        worksheet.column_dimensions['B'].width = 72
        buffer = BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()

    def _build_project_pptx(self, project):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = project.name
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, line in enumerate(self._project_document_lines(project)[1:]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = line
        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def _competition_document_lines(self, competition, version):
        event = competition.event
        revision_focus = (
            '首轮版本：围绕赛道要求完成结构搭建与证据清单。'
            if version == 1
            else (
                f'本届修订：结合 {event.edition} 年评审口径，'
                '更新数据、商业论证、成员分工与答辩重点。'
            )
        )
        return [
            f'项目名称：{competition.project.name}',
            f'项目编号：{competition.project.code}',
            f'项目负责人：{competition.project.leader.name}',
            f'比赛名称：{event.name}',
            f'比赛届次：{event.edition}',
            f'参赛队：{competition.entry_name}',
            f'参赛赛道：{competition.comp_type}',
            f'材料版本：V{version}',
            revision_focus,
        ]

    def _build_competition_plan_docx(self, competition, version):
        document = Document()
        document.add_heading(
            f'{competition.event.name} {competition.event.edition} 参赛计划书',
            level=1,
        )
        for line in self._competition_document_lines(competition, version):
            document.add_paragraph(line)
        document.add_heading('本届独立内容', level=2)
        document.add_paragraph(
            f'本计划书仅用于“{competition.entry_name}”，不得作为其他比赛、'
            '其他项目或其他年度的共用文件。'
        )
        buffer = BytesIO()
        document.save(buffer)
        return buffer.getvalue()

    def _build_competition_pitch_pptx(self, competition):
        presentation = Presentation()
        title_slide = presentation.slides.add_slide(
            presentation.slide_layouts[1]
        )
        title_slide.shapes.title.text = (
            f'{competition.project.name}｜{competition.event.edition}'
            f'{competition.event.name}'
        )
        body = title_slide.placeholders[1].text_frame
        body.clear()
        for index, line in enumerate(
            self._competition_document_lines(competition, 2)[1:]
        ):
            paragraph = (
                body.paragraphs[0]
                if index == 0
                else body.add_paragraph()
            )
            paragraph.text = line
        unique_slide = presentation.slides.add_slide(
            presentation.slide_layouts[1]
        )
        unique_slide.shapes.title.text = '本届路演重点'
        unique_slide.placeholders[1].text = (
            f'{competition.entry_name}围绕“{competition.comp_type}”'
            f'为{competition.event.edition}届单独准备，'
            '项目数据、现场展示和问答口径均与其他比赛材料区分。'
        )
        buffer = BytesIO()
        presentation.save(buffer)
        return buffer.getvalue()

    def _build_project_file(self, project, suffix):
        builders = {
            '.pdf': self._build_project_pdf,
            '.docx': self._build_project_docx,
            '.xlsx': self._build_project_xlsx,
            '.pptx': self._build_project_pptx,
        }
        try:
            return builders[suffix](project)
        except KeyError as exc:
            raise CommandError(f'不支持的演示文件类型：{suffix}') from exc

    def create_demo_files(self):
        """创建项目基础文件及每个参赛条目的独立计划书和路演 PPT。"""
        self.stdout.write('-> 创建项目文件、逐届比赛材料、版本与任务附件...')

        created_files = []
        version_count = 0
        for project in self.projects:
            project_files = []
            for label, suffix, content_type, level in FILE_SPECS:
                data = self._build_project_file(project, suffix)
                version = 2 if suffix == '.docx' else 1
                asset = FileAsset(
                    project=project,
                    name=f'{project.code} {label}{suffix}',
                    level=level,
                    size=len(data),
                    content_type=content_type,
                    uploader=project.leader,
                    version=version,
                    watermark_text=(
                        f'{DEMO_MARKER}{project.code}'
                        if level == FileAsset.Level.PUBLIC
                        else ''
                    ),
                )
                storage_name = (
                    f'{DEMO_IMPORT_DIRNAME}/{project.code}_'
                    f'{suffix.lstrip(".")}{suffix}'
                )
                asset.file.save(storage_name, ContentFile(data), save=True)
                asset_created_at = timezone.make_aware(
                    datetime.combine(
                        project.start_date + timedelta(days=30),
                        time(9, 0),
                    )
                )
                FileAsset.all_objects.filter(pk=asset.pk).update(
                    created_at=asset_created_at,
                    updated_at=asset_created_at,
                )
                project_files.append(asset)
                created_files.append(asset)

                if suffix == '.docx':
                    historical = FileVersion(
                        file_asset=asset,
                        version=1,
                        uploader=project.leader,
                    )
                    historical.file.save(
                        f'{DEMO_IMPORT_DIRNAME}/{project.code}_work_plan_v1.docx',
                        ContentFile(data),
                        save=True,
                    )
                    FileVersion.objects.filter(pk=historical.pk).update(
                        created_at=asset_created_at - timedelta(days=14)
                    )
                    version_count += 1

            self.files_by_project[project.id] = project_files
            tasks = self.tasks_by_project.get(project.id, [])
            if tasks:
                tasks[0].attachment_files.add(project_files[0], project_files[1])
            if len(tasks) > 1:
                tasks[1].attachment_files.add(project_files[2])
            if len(tasks) > 2:
                tasks[2].attachment_files.add(project_files[3])

        for competition in self.competitions:
            project = competition.project
            event = competition.event
            entry_key = f'entry_{competition.pk}'

            plan_v1 = self._build_competition_plan_docx(competition, 1)
            plan_v2 = self._build_competition_plan_docx(competition, 2)
            plan_asset = FileAsset(
                project=project,
                competition_entry=competition,
                name=(
                    f'{project.code}｜{event.edition}{event.name}｜'
                    '参赛计划书.docx'
                ),
                level=FileAsset.Level.INTERNAL,
                size=len(plan_v2),
                content_type=(
                    'application/vnd.openxmlformats-officedocument.'
                    'wordprocessingml.document'
                ),
                uploader=project.leader,
                version=2,
            )
            plan_asset.file.save(
                (
                    f'{DEMO_IMPORT_DIRNAME}/{entry_key}_'
                    f'{event.edition}_plan_v2.docx'
                ),
                ContentFile(plan_v2),
                save=True,
            )
            plan_history = FileVersion(
                file_asset=plan_asset,
                version=1,
                uploader=project.leader,
            )
            plan_history.file.save(
                (
                    f'{DEMO_IMPORT_DIRNAME}/{entry_key}_'
                    f'{event.edition}_plan_v1.docx'
                ),
                ContentFile(plan_v1),
                save=True,
            )
            material_history_at = timezone.make_aware(
                datetime.combine(
                    competition.register_date + timedelta(days=14),
                    time(9, 0),
                )
            )
            material_current_at = timezone.make_aware(
                datetime.combine(
                    competition.material_deadline - timedelta(days=5),
                    time(18, 0),
                )
            )
            FileVersion.objects.filter(pk=plan_history.pk).update(
                created_at=material_history_at
            )
            FileAsset.all_objects.filter(pk=plan_asset.pk).update(
                created_at=material_current_at,
                updated_at=material_current_at,
            )

            pitch_data = self._build_competition_pitch_pptx(competition)
            pitch_asset = FileAsset(
                project=project,
                competition_entry=competition,
                name=(
                    f'{project.code}｜{event.edition}{event.name}｜'
                    '路演PPT.pptx'
                ),
                level=FileAsset.Level.INTERNAL,
                size=len(pitch_data),
                content_type=(
                    'application/vnd.openxmlformats-officedocument.'
                    'presentationml.presentation'
                ),
                uploader=project.leader,
                version=1,
            )
            pitch_asset.file.save(
                (
                    f'{DEMO_IMPORT_DIRNAME}/{entry_key}_'
                    f'{event.edition}_pitch.pptx'
                ),
                ContentFile(pitch_data),
                save=True,
            )
            pitch_created_at = timezone.make_aware(
                datetime.combine(
                    competition.defense_date - timedelta(days=10),
                    time(18, 0),
                )
            )
            FileAsset.all_objects.filter(pk=pitch_asset.pk).update(
                created_at=pitch_created_at,
                updated_at=pitch_created_at,
            )

            entry_files = [plan_asset, pitch_asset]
            self.competition_files_by_entry[competition.id] = entry_files
            created_files.extend(entry_files)
            version_count += 1

        extension_counts = Counter(
            Path(asset.file.name).suffix.lower() for asset in created_files
        )
        self.stdout.write(self.style.SUCCESS(
            f'   文件创建完成：{len(created_files)} 个，历史版本 {version_count} 个，'
            f'PDF/DOCX/XLSX/PPTX = '
            f'{extension_counts[".pdf"]}/{extension_counts[".docx"]}/'
            f'{extension_counts[".xlsx"]}/{extension_counts[".pptx"]}'
        ))

    def create_import_history(self):
        """创建所有可导入模块的历史记录，并保留可读取的演示源文件。"""
        self.stdout.write('-> 创建导入历史...')
        source_files = sorted(self.assets_dir().rglob('*.xlsx'))
        if not source_files:
            raise CommandError('缺少 XLSX 演示模板，无法创建导入历史')

        import_dir = (
            Path(settings.MEDIA_ROOT)
            / 'imports'
            / DEMO_IMPORT_DIRNAME
        )
        import_dir.mkdir(parents=True, exist_ok=True)
        project = self.projects[0]
        member = self.members[0]
        definitions = [
            {
                'module': ImportTask.Module.PROJECTS,
                'status': ImportTask.Status.CONFIRMED,
                'mapping': {'项目名称': 'name', '项目编号': 'code', '负责人邮箱': 'leader_email'},
                'rows': [{'项目名称': project.name, '项目编号': project.code, '负责人邮箱': project.leader.email}],
            },
            {
                'module': ImportTask.Module.HISTORY_PROJECTS,
                'status': ImportTask.Status.ROLLED_BACK,
                'mapping': {'项目名称': 'name', '项目编号': 'code', '实际结束日期': 'actual_end_date'},
                'rows': [{'项目名称': self.projects[1].name, '项目编号': self.projects[1].code, '实际结束日期': str(self.projects[1].actual_end_date or '')}],
            },
            {
                'module': ImportTask.Module.MEMBERS,
                'status': ImportTask.Status.CONFIRMED,
                'mapping': {'姓名': 'name', '邮箱': 'email', '成员状态': 'membership_status'},
                'rows': [{'姓名': member.name, '邮箱': member.email, '成员状态': member.membership_status}],
            },
            {
                'module': ImportTask.Module.COMPETITIONS,
                'status': ImportTask.Status.PREVIEWED,
                'mapping': {'比赛名称': 'name', '项目编号': 'project_code', '级别': 'level'},
                'rows': [{'比赛名称': COMPETITION_DATA[0]['name'], '项目编号': project.code, '级别': 'national'}],
            },
            {
                'module': ImportTask.Module.TASKS,
                'status': ImportTask.Status.CONFIRMED,
                'mapping': {'任务标题': 'title', '项目编号': 'project_code', '负责人邮箱': 'assignee_email'},
                'rows': [{'任务标题': self.tasks_by_project[project.id][0].title, '项目编号': project.code, '负责人邮箱': member.email}],
            },
            {
                'module': ImportTask.Module.FINANCE,
                'status': ImportTask.Status.FAILED,
                'mapping': {'支出标题': 'title', '金额': 'amount', '项目编号': 'project_code', '支出日期': 'expense_date'},
                'rows': [{'支出标题': '调研交通费', '金额': '待确认', '项目编号': project.code, '支出日期': str(timezone.localdate())}],
                'errors': {'1': ['金额格式不正确，请填写数字。']},
            },
            {
                'module': ImportTask.Module.IP_APPLICATIONS,
                'status': ImportTask.Status.PREVIEWED,
                'mapping': {'成果名称': 'title', '内部编号': 'application_code', '关联项目编号': 'related_project_code'},
                'rows': [{'成果名称': '校园导览软件', '内部编号': f'{DEMO_IP_PREFIX}PREVIEW', '关联项目编号': project.code}],
            },
        ]

        for index, definition in enumerate(definitions):
            source = source_files[index % len(source_files)]
            destination = import_dir / f'{index + 1:02d}_{definition["module"]}.xlsx'
            destination.write_bytes(source.read_bytes())
            rows = definition['rows']
            errors = definition.get('errors', {})
            headers = list(rows[0].keys()) if rows else []
            import_task = ImportTask.objects.create(
                module=definition['module'],
                file_path=str(destination),
                status=definition['status'],
                field_mapping=definition['mapping'],
                preview_data={
                    'headers': headers,
                    'rows': rows,
                    'total_preview': len(rows),
                },
                snapshot=[],
                total_rows=len(rows),
                valid_rows=len(rows) - len(errors),
                error_rows=len(errors),
                error_details=errors,
                created_by=self.users['admin'],
            )
            ImportTask.objects.filter(pk=import_task.pk).update(
                created_at=timezone.now() - timedelta(days=21 - index * 3),
                updated_at=timezone.now() - timedelta(days=20 - index * 3),
            )

        material_path = import_dir / '08_materials.zip'
        material_manifest = {
            'version': 1,
            'items': [
                {
                    'path': '项目资料/协作说明.txt',
                    'name': '项目协作说明',
                    'project_code': project.code,
                    'level': 'internal',
                    'visibility': 'project',
                },
            ],
        }
        with zipfile.ZipFile(
            material_path,
            'w',
            compression=zipfile.ZIP_DEFLATED,
        ) as archive:
            archive.writestr(
                'manifest.json',
                json.dumps(material_manifest, ensure_ascii=False),
            )
            archive.writestr(
                '项目资料/协作说明.txt',
                '本资料包用于演示 ZIP + manifest.json 安全预览和确认导入。',
            )
        material_preview = preview_material_archive(
            material_path,
            team=self.team,
            operator=self.users['admin'],
        )
        material_task = ImportTask.objects.create(
            module=ImportTask.Module.MATERIALS,
            file_path=str(material_path),
            status=ImportTask.Status.PREVIEWED,
            field_mapping={},
            preview_data={
                'archive_sha256': material_preview['archive_sha256'],
                'rows': material_preview['rows'],
                'manifest_version': 1,
            },
            snapshot=[],
            total_rows=material_preview['total_rows'],
            valid_rows=material_preview['valid_rows'],
            error_rows=material_preview['error_rows'],
            error_details=material_preview['errors'],
            created_by=self.users['admin'],
            team=self.team,
        )
        ImportTask.objects.filter(pk=material_task.pk).update(
            created_at=timezone.now(),
            updated_at=timezone.now(),
        )

        self.stdout.write(self.style.SUCCESS(
            f'   导入历史创建完成：{len(definitions) + 1} 条，覆盖全部模块'
        ))

    # ------------------------------------------------------------------
    # 技能标签与成员技能
    # ------------------------------------------------------------------
    def create_skills(self):
        self.stdout.write('-> 创建技能标签与成员技能...')
        skills = []
        for name in SKILL_NAMES:
            skill, _ = SkillTag.objects.get_or_create(name=name)
            skills.append(skill)

        all_members = self.members + self.leaders
        # 为约 20 个成员分配 2~3 个技能
        chosen_members = random.sample(all_members, min(20, len(all_members)))
        count = 0
        for user in chosen_members:
            num = random.randint(2, 3)
            picked = random.sample(skills, num)
            for skill in picked:
                MemberSkill.objects.update_or_create(
                    user=user,
                    skill=skill,
                    defaults={'proficiency': random.randint(1, 5)},
                )
                count += 1

        self.stdout.write(self.style.SUCCESS(
            f'   技能创建完成：标签 {len(skills)} 个，成员技能 {count} 条'
        ))

    # ------------------------------------------------------------------
    # 灵活工作时间
    # ------------------------------------------------------------------
    def create_work_schedules(self):
        self.stdout.write('-> 创建灵活工作时间...')
        today = timezone.localdate()
        if today.day <= 15:
            period_start = today.replace(day=1)
            period_end = today.replace(day=15)
        else:
            period_start = today.replace(day=16)
            last_day = calendar.monthrange(today.year, today.month)[1]
            period_end = today.replace(day=last_day)

        all_members = self.members + self.leaders
        chosen = random.sample(all_members, min(15, len(all_members)))
        count = 0
        for user in chosen:
            work_hours = Decimal(random.randint(20, 80))
            detail = {
                '周一': random.randint(2, 8),
                '周二': random.randint(2, 8),
                '周三': random.randint(2, 8),
                '周四': random.randint(2, 8),
                '周五': random.randint(2, 8),
                '周末': random.randint(0, 6),
                '备注': '本周期可投入工时如上，课表已避开。',
            }
            FlexibleWorkSchedule.objects.update_or_create(
                user=user,
                period_start=period_start,
                defaults={
                    'period_end': period_end,
                    'work_hours': work_hours,
                    'detail': detail,
                    'can_offline': random.choice([True, False]),
                    'can_urgent': random.choice([True, False]),
                    'is_saturated': random.choice([True, False, False]),
                    'notes': (
                        f'{DEMO_MARKER}'
                        + random.choice([
                            '本周期安排稳定',
                            '近期有考试，工时偏少',
                            '可接受紧急任务',
                            '已满负荷',
                        ])
                    ),
                },
            )
            count += 1

        self.stdout.write(self.style.SUCCESS(f'   灵活工时创建完成：{count} 条'))

    # ------------------------------------------------------------------
    # 贡献记录
    # ------------------------------------------------------------------
    def create_contributions(self):
        self.stdout.write('-> 创建贡献记录...')
        now = timezone.now()
        period = now.strftime('%Y-%m')

        # 15 条：8 approved / 5 pending / 2 rejected
        statuses = (
            [Contribution.Status.APPROVED] * 8 +
            [Contribution.Status.PENDING] * 5 +
            [Contribution.Status.REJECTED] * 2
        )
        random.shuffle(statuses)

        contrib_types = [
            Contribution.ContributionType.STAGE_TASK,
            Contribution.ContributionType.CORE,
            Contribution.ContributionType.LONG_TERM,
            Contribution.ContributionType.RESOURCE,
            Contribution.ContributionType.TEMPORARY_HELP,
            Contribution.ContributionType.COMPETITION,
            Contribution.ContributionType.FINANCE_MANAGE,
            Contribution.ContributionType.IP_WRITING,
        ]

        contents = [
            '负责项目核心模块开发，按时完成阶段性目标。',
            '承担大量文档撰写与材料整理工作。',
            '长期参与项目，持续贡献代码与思路。',
            '为团队提供关键测试设备与数据资源。',
            '临时协助完成答辩PPT制作。',
            '作为核心成员主导技术方案设计。',
            '负责比赛申报材料撰写并顺利提交。',
            '协助完成经费报销与票据整理。',
        ]

        for i, status in enumerate(statuses):
            project = self.projects[i % len(self.projects)]
            members = self.project_members[project.id]
            user = random.choice(members)
            leader = project.leader
            ctype = random.choice(contrib_types)
            contribution_summary = random.choice(contents)
            content = (
                f'{user.name}在「{project.name}」中'
                f'{contribution_summary}'
            )
            weight = Decimal(random.randint(1, 10))

            kwargs = dict(
                user=user,
                project=project,
                contribution_type=ctype,
                content=content,
                description=content,
                weight=weight,
                status=status,
                period=period,
                filled_by=leader,
            )
            if status == Contribution.Status.APPROVED:
                kwargs.update(
                    reviewer=leader,
                    reviewed_at=now - timedelta(days=random.randint(1, 10)),
                    review_opinion='贡献属实，予以确认。',
                    score=weight * Decimal('2'),
                )
            elif status == Contribution.Status.REJECTED:
                kwargs.update(
                    reviewer=leader,
                    reviewed_at=now - timedelta(days=random.randint(1, 5)),
                    review_opinion='描述与实际不符，暂不确认。',
                )
            Contribution.objects.create(**kwargs)

        self.stdout.write(self.style.SUCCESS(
            f'   贡献记录创建完成：{len(statuses)} 条'
        ))

    # ------------------------------------------------------------------
    # 成员排序与异议
    # ------------------------------------------------------------------
    def create_rankings(self):
        self.stdout.write('-> 创建成员排序与异议...')
        now = timezone.now()
        period = now.strftime('%Y-%m')

        # 前 2 个项目：第 1 个 draft，第 2 个 confirmed
        ranking_specs = [
            (self.projects[0], MemberRanking.Status.DRAFT),
            (self.projects[1], MemberRanking.Status.CONFIRMED),
        ]

        all_rankings = []
        for project, status in ranking_specs:
            members = self.project_members[project.id]
            ranked_members = random.sample(members, min(5, len(members)))
            # 按随机的“分值”降序排名
            scored = [(m, Decimal(random.randint(10, 100))) for m in ranked_members]
            scored.sort(key=lambda x: x[1], reverse=True)

            is_public = status == MemberRanking.Status.CONFIRMED
            for rank, (user, score) in enumerate(scored, start=1):
                ranking = MemberRanking.objects.create(
                    user=user,
                    project=project,
                    period=period,
                    status=status,
                    total_score=score,
                    rank=rank,
                    task_completed_count=random.randint(0, 12),
                    project_count=random.randint(1, 3),
                    competition_count=random.randint(0, 2),
                    ip_contribution_count=random.randint(0, 2),
                    is_published=is_public,
                    is_public=is_public,
                )
                all_rankings.append(ranking)

        # 2 条异议：1 pending / 1 leader_reviewed
        if len(all_rankings) >= 2:
            # 第一条：pending（来自 draft 项目，提出人为该项目成员）
            r1 = all_rankings[0]
            objector1 = random.choice(self.project_members[r1.project_id])
            RankingObjection.objects.create(
                ranking=r1,
                objector=objector1,
                content='我对自己的排名有异议，本周期完成了较多阶段性任务，'
                        '认为贡献分值应更高。',
                status=RankingObjection.Status.PENDING,
            )
            # 第二条：leader_reviewed（来自 confirmed 项目）
            r2 = all_rankings[-1]
            members2 = self.project_members[r2.project_id]
            objector2 = random.choice(members2)
            RankingObjection.objects.create(
                ranking=r2,
                objector=objector2,
                content='排名结果未充分考虑 IP 撰写贡献，请求复核。',
                status=RankingObjection.Status.LEADER_REVIEWED,
                leader_opinion='经核实，该成员 IP 贡献确有遗漏，建议老师最终确认后调整。',
                leader_reviewer=self.projects[1].leader,
                leader_reviewed_at=now - timedelta(days=1),
            )

        self.stdout.write(self.style.SUCCESS(
            f'   排序创建完成：{len(all_rankings)} 条，异议 2 条'
        ))

    # ------------------------------------------------------------------
    # 定时报表
    # ------------------------------------------------------------------
    def create_scheduled_reports(self):
        """创建可直接下载、也可继续按计划运行的三种格式演示报表。"""
        self.stdout.write('-> 创建定时报表与成功执行记录...')
        now = timezone.now()
        definitions = [
            {
                'name': '每日项目执行概览',
                'description': '每天汇总进行中、暂停及已结项项目。',
                'report_type': CustomReport.ReportType.SUMMARY,
                'config': {
                    'data_source': 'project',
                    'group_by': 'status',
                    'chart_type': 'table',
                    'filters': {},
                },
                'creator': self.users['admin'],
                'recipients': [
                    self.users['admin'],
                    self.users['teacher1'],
                ],
                'frequency': ScheduledReport.Frequency.DAILY,
                'execution_time': time(8, 30),
                'file_format': ScheduledReport.FileFormat.XLSX,
                'project': self.projects[0],
                'schedule_cron': '30 8 * * *',
            },
            {
                'name': '每周任务交付简报',
                'description': '每周汇总任务完成、审核及延期情况。',
                'report_type': CustomReport.ReportType.TREND,
                'config': {
                    'data_source': 'task',
                    'group_by': 'project',
                    'chart_type': 'table',
                    'filters': {},
                },
                'creator': self.users['teacher1'],
                'recipients': [
                    self.users['teacher1'],
                    self.users['leader1'],
                ],
                'frequency': ScheduledReport.Frequency.WEEKLY,
                'execution_time': time(9, 0),
                'weekday': 0,
                'file_format': ScheduledReport.FileFormat.DOCX,
                'project': self.projects[1],
                'schedule_cron': '0 9 * * 1',
            },
            {
                'name': '每月经费执行报告',
                'description': '每月汇总项目支出和经费使用结构。',
                'report_type': CustomReport.ReportType.COMPARISON,
                'config': {
                    'data_source': 'finance',
                    'group_by': 'project',
                    'chart_type': 'table',
                    'filters': {},
                },
                'creator': self.users['leader1'],
                'recipients': [
                    self.users['leader1'],
                    self.users['teacher2'],
                ],
                'frequency': ScheduledReport.Frequency.MONTHLY,
                'execution_time': time(10, 0),
                'day_of_month': 5,
                'file_format': ScheduledReport.FileFormat.PDF,
                'project': self.projects[2],
                'schedule_cron': '0 10 5 * *',
            },
        ]
        builders = {
            ScheduledReport.FileFormat.XLSX: self._build_project_xlsx,
            ScheduledReport.FileFormat.DOCX: self._build_project_docx,
            ScheduledReport.FileFormat.PDF: self._build_project_pdf,
        }

        for index, definition in enumerate(definitions, start=1):
            report = CustomReport.objects.create(
                name=f'{DEMO_MARKER}{definition["name"]}',
                description=definition['description'],
                report_type=definition['report_type'],
                config=definition['config'],
                created_by=definition['creator'],
                is_scheduled=True,
                schedule_cron=definition['schedule_cron'],
            )
            schedule = ScheduledReport.objects.create(
                report=report,
                created_by=definition['creator'],
                frequency=definition['frequency'],
                execution_time=definition['execution_time'],
                weekday=definition.get('weekday', 0),
                day_of_month=definition.get('day_of_month', 1),
                timezone='Asia/Shanghai',
                file_format=definition['file_format'],
                last_run=now,
                last_status=ScheduledReport.RunStatus.SUCCESS,
                is_active=True,
            )
            schedule.recipients.set(definition['recipients'])
            schedule.next_run = compute_next_run(schedule, base=now)
            schedule.save(update_fields=['next_run'])

            content = builders[definition['file_format']](
                definition['project']
            )
            file_name = (
                f'{definition["name"]}.{definition["file_format"]}'
            )
            execution = ScheduledReportExecution(
                schedule=schedule,
                trigger=ScheduledReportExecution.Trigger.SCHEDULED,
                status=ScheduledReport.RunStatus.SUCCESS,
                file_name=file_name,
                file_format=definition['file_format'],
                file_size=len(content),
                delivery_status=(
                    ScheduledReportExecution.DeliveryStatus.NOT_REQUESTED
                ),
                recipient_snapshot=[
                    {
                        'id': recipient.pk,
                        'name': recipient.name,
                        'email': recipient.email,
                    }
                    for recipient in definition['recipients']
                ],
                message='演示报表已生成并保留站内文件，未实际发送邮件。',
                finished_at=now,
                generated_by=definition['creator'],
            )
            execution.file.save(
                (
                    f'{DEMO_IMPORT_DIRNAME}/scheduled_report_{index}.'
                    f'{definition["file_format"]}'
                ),
                ContentFile(content),
                save=False,
            )
            execution.save()

        self.stdout.write(self.style.SUCCESS(
            '   定时报表创建完成：3 个计划、3 个成功执行文件'
        ))

    # ------------------------------------------------------------------
    # 知识产权申请
    # ------------------------------------------------------------------
    def _build_ip_certificate_pdf(
        self,
        project,
        application_code,
        title,
        authorized_date,
    ):
        """生成可核验、内容明确的小型演示授权登记证书 PDF。"""
        buffer = BytesIO()
        document = canvas.Canvas(buffer, pagesize=A4)
        document.setTitle(f'{application_code} 授权登记证书')
        document.setAuthor('团队管理软件演示数据')
        font_name = self._pdf_font_name()
        document.setFont(font_name, 20)
        document.drawCentredString(297, 775, '知识产权授权登记证书（演示）')
        document.setFont(font_name, 12)
        lines = (
            f'申请编号：{application_code}',
            f'成果名称：{title}',
            f'关联项目：{project.code} {project.name}',
            f'项目负责人：{project.leader.name}',
            f'授权登记日期：{authorized_date.isoformat()}',
            '证书状态：授权登记完成，可进入成果归档流程。',
            '说明：本证书由 seed_demo_data 生成，仅用于系统功能演示。',
        )
        y = 710
        for line in lines:
            document.drawString(72, y, line)
            y -= 38
        document.rect(54, 440, 487, 320)
        document.save()
        return buffer.getvalue()

    def _create_ip_certificate_asset(
        self,
        project,
        application_code,
        title,
        authorized_date,
    ):
        """创建独立内部证书文件，避免把普通项目文档冒充授权证书。"""
        content = self._build_ip_certificate_pdf(
            project,
            application_code,
            title,
            authorized_date,
        )
        certificate = FileAsset(
            project=project,
            name=f'{application_code} 最终授权登记证书.pdf',
            level=FileAsset.Level.INTERNAL,
            size=len(content),
            content_type='application/pdf',
            uploader=project.leader,
        )
        certificate.file.save(
            f'{DEMO_IMPORT_DIRNAME}/{application_code}_certificate.pdf',
            ContentFile(content),
            save=True,
        )
        return certificate

    def _create_legacy_ip_applications(self):
        self.stdout.write('-> 创建知识产权申请...')
        now = timezone.now()

        def pick_members(project, n):
            members = self.project_members[project.id]
            return random.sample(members, min(n, len(members)))

        # --- 1. 智能校园导览系统 软著 - 已授权 ---
        p1 = self.projects[0]
        writer1 = p1.leader
        ip1_code = f'{DEMO_IP_PREFIX}001'
        ip1_title = '智能校园导览系统'
        ip1_authorized_date = (now - timedelta(days=30)).date()
        certificate_file = self._create_ip_certificate_asset(
            p1,
            ip1_code,
            ip1_title,
            ip1_authorized_date,
        )
        ip1 = IntellectualPropertyApplication.objects.create(
            title=ip1_title,
            application_code=ip1_code,
            ip_type=IntellectualPropertyApplication.IPType.SOFTWARE_COPYRIGHT,
            related_project=p1,
            status=IntellectualPropertyApplication.Status.AUTHORIZED,
            main_writer=writer1,
            applicant_executor=writer1,
            material_manager=pick_members(p1, 1)[0],
            project_reviewer=p1.leader,
            teacher_confirmer=self.teachers[0],
            start_date=(now - timedelta(days=180)).date(),
            submit_date=(now - timedelta(days=150)).date(),
            accepted_date=(now - timedelta(days=90)).date(),
            authorized_date=ip1_authorized_date,
            return_count=0,
            final_certificate_file=certificate_file,
            intro='基于移动端定位与AR技术的校园导览系统软件著作权，'
                  '包含室内外导航、景点讲解与活动指引等核心模块。',
            created_by=writer1,
        )
        self._add_ip_contributors(ip1, writer1, p1, [
            IPApplicationContributor.ContributorRole.MAIN_WRITER,
            IPApplicationContributor.ContributorRole.EXECUTOR,
            IPApplicationContributor.ContributorRole.MATERIAL_MANAGER,
            IPApplicationContributor.ContributorRole.REVIEWER,
        ])

        # --- 2. 医学影像辅助诊断算法 发明专利 - 科研处退回 ---
        p2 = self.projects[1]
        writer2 = p2.leader
        ip2 = IntellectualPropertyApplication.objects.create(
            title='医学影像辅助诊断算法',
            application_code=f'{DEMO_IP_PREFIX}002',
            ip_type=IntellectualPropertyApplication.IPType.INVENTION_PATENT,
            related_project=p2,
            status=IntellectualPropertyApplication.Status.RETURNED,
            main_writer=writer2,
            applicant_executor=writer2,
            material_manager=pick_members(p2, 1)[0],
            project_reviewer=p2.leader,
            teacher_confirmer=self.teachers[0],
            start_date=(now - timedelta(days=120)).date(),
            submit_date=(now - timedelta(days=60)).date(),
            return_count=1,
            current_problem='科研处退回：权利要求书保护范围不够清晰，说明书技术方案描述需补充实施例。',
            intro='一种基于深度卷积神经网络的医学影像病灶检测与分类方法及系统。',
            created_by=writer2,
        )
        self._add_ip_contributors(ip2, writer2, p2, [
            IPApplicationContributor.ContributorRole.MAIN_WRITER,
            IPApplicationContributor.ContributorRole.CODE_PROVIDER,
            IPApplicationContributor.ContributorRole.EXECUTOR,
            IPApplicationContributor.ContributorRole.REVIEWER,
        ])
        # 退回记录
        IPReturnRecord.objects.create(
            application=ip2,
            return_time=now - timedelta(days=7),
            return_source=IPReturnRecord.ReturnSource.RESEARCH_OFFICE,
            return_reason='权利要求保护范围过宽，说明书缺少足够实施例支撑，需补充技术效果描述。',
            responsibility_type=IPReturnRecord.ResponsibilityType.WRITING_PROBLEM,
            responsible_user=writer2,
            assigned_by=self.teachers[0],
            modify_deadline=now + timedelta(days=7),
            actual_modifier=writer2,
            modify_description='',
            result=IPReturnRecord.ReturnResult.PENDING,
        )

        # --- 3. 农产品溯源链码生成方法 发明专利 - 科研处审核中 ---
        p4 = self.projects[3]
        writer3 = p4.leader
        ip3 = IntellectualPropertyApplication.objects.create(
            title='农产品溯源链码生成方法',
            application_code=f'{DEMO_IP_PREFIX}003',
            ip_type=IntellectualPropertyApplication.IPType.INVENTION_PATENT,
            related_project=p4,
            status=IntellectualPropertyApplication.Status.RESEARCH_OFFICE_REVIEW,
            main_writer=writer3,
            applicant_executor=writer3,
            material_manager=pick_members(p4, 1)[0],
            project_reviewer=p4.leader,
            teacher_confirmer=self.teachers[0],
            start_date=(now - timedelta(days=90)).date(),
            submit_date=(now - timedelta(days=30)).date(),
            return_count=0,
            intro='一种基于联盟链的农产品溯源链码生成与校验方法，确保数据不可篡改。',
            created_by=writer3,
        )
        self._add_ip_contributors(ip3, writer3, p4, [
            IPApplicationContributor.ContributorRole.MAIN_WRITER,
            IPApplicationContributor.ContributorRole.DRAWING_PROVIDER,
            IPApplicationContributor.ContributorRole.EXECUTOR,
            IPApplicationContributor.ContributorRole.REVIEWER,
        ])

        # --- 4. 校园二手交易小程序 软著 - 材料撰写中 ---
        p3 = self.projects[2]
        writer4 = p3.leader
        ip4 = IntellectualPropertyApplication.objects.create(
            title='校园二手交易小程序',
            application_code=f'{DEMO_IP_PREFIX}004',
            ip_type=IntellectualPropertyApplication.IPType.SOFTWARE_COPYRIGHT,
            related_project=p3,
            status=IntellectualPropertyApplication.Status.WRITING,
            main_writer=writer4,
            applicant_executor=writer4,
            material_manager=pick_members(p3, 1)[0],
            project_reviewer=p3.leader,
            teacher_confirmer=self.teachers[0],
            start_date=(now - timedelta(days=30)).date(),
            return_count=0,
            current_problem='软件说明书撰写中，源代码文档待整理。',
            intro='面向在校学生的二手物品交易小程序软件著作权。',
            created_by=writer4,
        )
        self._add_ip_contributors(ip4, writer4, p3, [
            IPApplicationContributor.ContributorRole.MAIN_WRITER,
            IPApplicationContributor.ContributorRole.DOCUMENT_WRITER,
            IPApplicationContributor.ContributorRole.REVIEWER,
        ])

        # --- 5. 环境监测数据异常检测方法 实用新型专利 - 科研处退回 ---
        p9 = self.projects[8]
        writer5 = p9.leader
        ip5 = IntellectualPropertyApplication.objects.create(
            title='环境监测数据异常检测方法',
            application_code=f'{DEMO_IP_PREFIX}005',
            ip_type=IntellectualPropertyApplication.IPType.UTILITY_MODEL,
            related_project=p9,
            status=IntellectualPropertyApplication.Status.RETURNED,
            main_writer=writer5,
            applicant_executor=writer5,
            material_manager=pick_members(p9, 1)[0],
            project_reviewer=p9.leader,
            teacher_confirmer=self.teachers[0],
            start_date=(now - timedelta(days=150)).date(),
            submit_date=(now - timedelta(days=90)).date(),
            return_count=2,
            current_problem='第二次退回：附图标记与说明书不一致，权利要求项引用关系有误。',
            intro='一种用于环境监测物联网数据的异常检测与告警方法及装置。',
            created_by=writer5,
        )
        self._add_ip_contributors(ip5, writer5, p9, [
            IPApplicationContributor.ContributorRole.MAIN_WRITER,
            IPApplicationContributor.ContributorRole.DRAWING_PROVIDER,
            IPApplicationContributor.ContributorRole.MATERIAL_MANAGER,
            IPApplicationContributor.ContributorRole.REVIEWER,
        ])
        # 退回记录
        IPReturnRecord.objects.create(
            application=ip5,
            return_time=now - timedelta(days=3),
            return_source=IPReturnRecord.ReturnSource.RESEARCH_OFFICE,
            return_reason='附图标记与说明书描述不一致，权利要求项引用关系存在错误，需重新核对。',
            responsibility_type=IPReturnRecord.ResponsibilityType.MATERIAL_PROBLEM,
            responsible_user=pick_members(p9, 1)[0],
            assigned_by=self.teachers[0],
            modify_deadline=now + timedelta(days=10),
            actual_modifier=None,
            modify_description='',
            result=IPReturnRecord.ReturnResult.PENDING,
        )

        # 2 条知识产权异议
        IPObjection.objects.create(
            application=ip2,
            objector=pick_members(p2, 1)[0],
            objection_type=IPObjection.ObjectionType.RETURN_RESPONSIBILITY,
            content='退回责任不应全部由撰写人承担，材料整理环节也存在疏漏。',
            status=IPObjection.ObjectionStatus.PENDING,
        )
        IPObjection.objects.create(
            application=ip1,
            objector=pick_members(p1, 1)[0],
            objection_type=IPObjection.ObjectionType.WRITING_CREDIT,
            content='申请执行贡献认定偏低，请求复核贡献排序。',
            status=IPObjection.ObjectionStatus.RESOLVED,
            leader_opinion='经核实贡献属实，已调整排序。',
            leader_reviewer=p1.leader,
            leader_reviewed_at=now - timedelta(days=5),
            teacher_confirmer=self.teachers[0],
            teacher_confirmed_at=now - timedelta(days=3),
            final_result='维持原贡献认定，排序微调。',
        )

        material_specs = [
            (
                ip1,
                self.files_by_project[p1.id][0],
                IPMaterialVersion.MaterialType.ARCHIVE,
                'v1',
                True,
                '成果归档材料已核验，并与最终授权登记证书分别留存。',
            ),
            (
                ip2,
                self.files_by_project[p2.id][1],
                IPMaterialVersion.MaterialType.SPECIFICATION,
                'v2',
                False,
                '根据科研处退回意见补充实施例与技术效果。',
            ),
            (
                ip3,
                self.files_by_project[p4.id][1],
                IPMaterialVersion.MaterialType.DISCLOSURE,
                'v1',
                False,
                '科研处审核使用的技术交底书。',
            ),
            (
                ip4,
                self.files_by_project[p3.id][1],
                IPMaterialVersion.MaterialType.MANUAL,
                'v1',
                False,
                '软件说明书初稿，等待代码文档补充。',
            ),
            (
                ip5,
                self.files_by_project[p9.id][1],
                IPMaterialVersion.MaterialType.DRAWING,
                'v2',
                False,
                '已重新核对附图标记，待责任人确认。',
            ),
        ]
        for application, file_asset, material_type, version, is_final, note in material_specs:
            IPMaterialVersion.objects.create(
                application=application,
                file_asset=file_asset,
                material_type=material_type,
                version=version,
                uploaded_by=application.material_manager or application.main_writer,
                change_note=note,
                is_final=is_final,
            )

        self.stdout.write(self.style.SUCCESS(
            '   知识产权创建完成：申请 5 个，材料版本 5 条，退回记录 2 条，异议 2 条'
        ))

    def create_ip_applications(self):
        """创建 40 个专利、25 个软著和 7 个科技查新档案。"""
        self.stdout.write('-> 创建知识产权与科技查新档案...')
        now = timezone.now()
        application_specs = [
            *[
                (
                    IntellectualPropertyApplication.IPType.INVENTION_PATENT,
                    f'发明专利 {index + 1:02d}',
                )
                for index in range(22)
            ],
            *[
                (
                    IntellectualPropertyApplication.IPType.UTILITY_MODEL,
                    f'实用新型专利 {index + 1:02d}',
                )
                for index in range(14)
            ],
            *[
                (
                    IntellectualPropertyApplication.IPType.DESIGN_PATENT,
                    f'外观设计专利 {index + 1:02d}',
                )
                for index in range(4)
            ],
            *[
                (
                    IntellectualPropertyApplication.IPType.SOFTWARE_COPYRIGHT,
                    f'软件著作权 {index + 1:02d}',
                )
                for index in range(25)
            ],
            *[
                (
                    IntellectualPropertyApplication.IPType.NOVELTY_SEARCH,
                    f'科技查新报告 {index + 1:02d}',
                )
                for index in range(7)
            ],
        ]
        archived_statuses = (
            IntellectualPropertyApplication.Status.AUTHORIZED,
            IntellectualPropertyApplication.Status.ARCHIVED,
        )
        applications = []
        returned_applications = []

        for index, (ip_type, type_label) in enumerate(application_specs):
            project = self.projects[index % len(self.projects)]
            project_users = self.project_members[project.id]
            main_writer = project.leader
            applicant_executor = project_users[
                1 + index % (len(project_users) - 1)
            ]
            material_manager = project_users[
                1 + (index + 2) % (len(project_users) - 1)
            ]
            year = 2021 + index % 6
            start_date = date(year, 1 + index % 5, 5 + index % 18)

            if year <= 2023:
                status = archived_statuses[index % len(archived_statuses)]
            elif year == 2024:
                status = (
                    IntellectualPropertyApplication.Status.ACCEPTED,
                    IntellectualPropertyApplication.Status.AUTHORIZED,
                    IntellectualPropertyApplication.Status.ARCHIVED,
                )[index % 3]
            elif year == 2025:
                status = (
                    IntellectualPropertyApplication.Status.RESEARCH_OFFICE_REVIEW,
                    IntellectualPropertyApplication.Status.RETURNED,
                    IntellectualPropertyApplication.Status.ACCEPTED,
                    IntellectualPropertyApplication.Status.RESUBMITTED,
                )[index % 4]
            else:
                status = (
                    IntellectualPropertyApplication.Status.WRITING,
                    IntellectualPropertyApplication.Status.LEADER_REVIEW,
                    IntellectualPropertyApplication.Status.TEACHER_CONFIRM,
                    IntellectualPropertyApplication.Status.DRAFT,
                )[index % 4]

            submit_date = (
                start_date + timedelta(days=45)
                if status
                not in {
                    IntellectualPropertyApplication.Status.DRAFT,
                    IntellectualPropertyApplication.Status.WRITING,
                }
                else None
            )
            accepted_date = (
                start_date + timedelta(days=100)
                if status
                in {
                    IntellectualPropertyApplication.Status.ACCEPTED,
                    *archived_statuses,
                }
                else None
            )
            authorized_date = (
                start_date + timedelta(days=180)
                if status in archived_statuses
                else None
            )
            code = f'{DEMO_IP_PREFIX}{index + 1:03d}'
            title = f'{project.name}｜{type_label}'
            certificate_file = None
            if index == 0:
                status = IntellectualPropertyApplication.Status.AUTHORIZED
                submit_date = start_date + timedelta(days=45)
                accepted_date = start_date + timedelta(days=100)
                authorized_date = start_date + timedelta(days=180)
                certificate_file = self._create_ip_certificate_asset(
                    project,
                    code,
                    title,
                    authorized_date,
                )

            application = IntellectualPropertyApplication.objects.create(
                title=title,
                application_code=code,
                ip_type=ip_type,
                related_project=project,
                status=status,
                main_writer=main_writer,
                applicant_executor=applicant_executor,
                material_manager=material_manager,
                project_reviewer=project.leader,
                teacher_confirmer=self.teachers[0],
                start_date=start_date,
                submit_date=submit_date,
                accepted_date=accepted_date,
                authorized_date=authorized_date,
                return_count=1 if status == (
                    IntellectualPropertyApplication.Status.RETURNED
                ) else 0,
                current_problem=(
                    '学校系统退回：需核对材料字段、署名顺序与附件版本。'
                    if status == IntellectualPropertyApplication.Status.RETURNED
                    else ''
                ),
                status_note=(
                    f'{year} 年启动，按项目成果实际进度形成演示档案。'
                ),
                final_certificate_file=certificate_file,
                intro=(
                    f'{project.name}形成的{type_label}，'
                    '完整记录撰写、复核、提交、退回与归档责任。'
                ),
                created_by=main_writer,
            )
            created_at = timezone.make_aware(
                datetime.combine(start_date, time(9, 0))
            )
            IntellectualPropertyApplication.objects.filter(
                pk=application.pk
            ).update(created_at=created_at, updated_at=created_at)
            applications.append(application)
            if status == IntellectualPropertyApplication.Status.RETURNED:
                returned_applications.append(application)

            self._add_ip_contributors(
                application,
                main_writer,
                project,
                [
                    IPApplicationContributor.ContributorRole.MAIN_WRITER,
                    IPApplicationContributor.ContributorRole.EXECUTOR,
                    IPApplicationContributor.ContributorRole.MATERIAL_MANAGER,
                    IPApplicationContributor.ContributorRole.REVIEWER,
                ],
            )

            material_type = {
                IntellectualPropertyApplication.IPType.SOFTWARE_COPYRIGHT:
                    IPMaterialVersion.MaterialType.MANUAL,
                IntellectualPropertyApplication.IPType.NOVELTY_SEARCH:
                    IPMaterialVersion.MaterialType.ARCHIVE,
                IntellectualPropertyApplication.IPType.DESIGN_PATENT:
                    IPMaterialVersion.MaterialType.DRAWING,
                IntellectualPropertyApplication.IPType.UTILITY_MODEL:
                    IPMaterialVersion.MaterialType.SPECIFICATION,
                IntellectualPropertyApplication.IPType.INVENTION_PATENT:
                    IPMaterialVersion.MaterialType.DISCLOSURE,
            }[ip_type]
            if index == 0:
                material_type = IPMaterialVersion.MaterialType.ARCHIVE
            IPMaterialVersion.objects.create(
                application=application,
                file_asset=self.files_by_project[project.id][
                    index % len(self.files_by_project[project.id])
                ],
                material_type=material_type,
                version='v2' if year <= 2025 else 'v1',
                uploaded_by=material_manager,
                change_note=(
                    f'{year} 年材料版本，已核对关联项目、责任人和成果类型。'
                ),
                is_final=status in archived_statuses,
            )

        for return_index, application in enumerate(returned_applications[:2]):
            IPReturnRecord.objects.create(
                application=application,
                return_time=now - timedelta(days=7 - return_index),
                return_source=IPReturnRecord.ReturnSource.RESEARCH_OFFICE,
                return_reason='材料字段与附件版本不一致，需完成核对后重新提交。',
                responsibility_type=(
                    IPReturnRecord.ResponsibilityType.MATERIAL_PROBLEM
                ),
                responsible_user=application.material_manager,
                assigned_by=self.teachers[0],
                modify_deadline=now + timedelta(days=7 + return_index),
                actual_modifier=application.material_manager,
                modify_description='按退回清单逐项核对中。',
                result=IPReturnRecord.ReturnResult.PENDING,
            )

        if returned_applications:
            IPObjection.objects.create(
                application=returned_applications[0],
                objector=returned_applications[0].material_manager,
                objection_type=IPObjection.ObjectionType.RETURN_RESPONSIBILITY,
                content='退回原因同时涉及撰写和系统字段，请负责人复核责任划分。',
                status=IPObjection.ObjectionStatus.PENDING,
            )
        IPObjection.objects.create(
            application=applications[0],
            objector=applications[0].applicant_executor,
            objection_type=IPObjection.ObjectionType.WRITING_CREDIT,
            content='申请执行与材料复核投入尚未在署名说明中充分体现。',
            status=IPObjection.ObjectionStatus.RESOLVED,
            leader_opinion='经核验贡献属实，已补充责任说明。',
            leader_reviewer=applications[0].project_reviewer,
            leader_reviewed_at=now - timedelta(days=5),
            teacher_confirmer=self.teachers[0],
            teacher_confirmed_at=now - timedelta(days=3),
            final_result='贡献说明已补充，原署名顺序保持不变。',
        )

        self.ip_apps = applications
        self.stdout.write(self.style.SUCCESS(
            '   成果创建完成：专利 40 个、软著 25 个、科技查新 7 个；'
            f'材料版本 {len(applications)} 条'
        ))

    def _add_ip_contributors(self, application, main_writer, project, roles):
        """按档案中的实际职责为 IP 申请建立一致、可确认的责任分工。"""
        members = [u for u in self.project_members[project.id] if u != main_writer]
        random.shuffle(members)
        direct_responsibilities = {
            IPApplicationContributor.ContributorRole.MAIN_WRITER: application.main_writer,
            IPApplicationContributor.ContributorRole.EXECUTOR: application.applicant_executor,
            IPApplicationContributor.ContributorRole.MATERIAL_MANAGER: application.material_manager,
            IPApplicationContributor.ContributorRole.REVIEWER: application.project_reviewer,
        }
        application_contributors = []
        for role in roles:
            user = direct_responsibilities.get(role)
            if user is None:
                user = members.pop() if members else main_writer
            application_contributors.append((user, role))
        for user, role in application_contributors:
            role_label = IPApplicationContributor.ContributorRole(role).label
            is_confirmed = (
                application.status
                in {
                    IntellectualPropertyApplication.Status.AUTHORIZED,
                    IntellectualPropertyApplication.Status.ARCHIVED,
                }
            )
            IPApplicationContributor.objects.create(
                application=application,
                user=user,
                role=role,
                contribution_description=(
                    f'{user.name}在「{application.title}」中担任'
                    f'{role_label}。'
                ),
                responsibility_description=(
                    f'{user.name}负责「{application.title}」对应环节的'
                    '材料质量、节点反馈与留痕，'
                    '如发生退回需参与原因说明和修订闭环。'
                ),
                is_confirmed=is_confirmed,
                confirmed_by=application.project_reviewer if is_confirmed else None,
                confirmed_at=timezone.now() - timedelta(days=25) if is_confirmed else None,
            )

    def create_portal_publications(self):
        """为演示项目、成果和已授权成员建立逐项公开决策。"""
        self.stdout.write('-> 创建受治理的公开门户内容...')
        teacher = self.users['teacher1']
        portal_defaults = {
            'team_name': '数智创新实践团队',
            'tagline': '真实项目 · 持续协作 · 成果沉淀',
            'summary': '团队自 2021 年起持续开展跨专业项目实践，所有公开内容均经过逐项确认。',
            'about_title': '让每一段协作过程都可追踪、可复盘',
            'about_text': '我们围绕校园、社区、乡村和产业真实需求开展项目，并沉淀任务、文件、赛事与知识产权成果。',
            'contact_email': 'teacher1@demo.com',
            'join_title': '加入数智创新实践团队',
            'join_message': '欢迎愿意持续投入、尊重协作规范并主动承担任务的同学联系我们。',
            'join_url': '/join-us',
            'updated_by': teacher,
        }
        portal_settings, created = PortalSettings.objects.get_or_create(
            singleton_key='default',
            defaults=portal_defaults,
        )
        if (
            not created
            and portal_settings.updated_by_id
            and portal_settings.updated_by.email.endswith('@demo.com')
        ):
            for field, value in portal_defaults.items():
                setattr(portal_settings, field, value)
            portal_settings.save(update_fields=[
                *portal_defaults.keys(),
                'updated_at',
            ])

        for index, project in enumerate(self.projects):
            PortalPublication.objects.update_or_create(
                content_type=PortalPublication.ContentType.PROJECT,
                object_id=project.id,
                defaults={
                    'is_public': index < 16,
                    'is_featured': index < 4,
                    'member_consent': False,
                    'display_order': index,
                    'custom_title': project.name,
                    'custom_summary': project.intro,
                    'updated_by': teacher,
                },
            )

        ip_applications = list(
            IntellectualPropertyApplication.objects.filter(
                application_code__startswith=DEMO_IP_PREFIX
            ).order_by('id')
        )
        for index, application in enumerate(ip_applications):
            PortalPublication.objects.update_or_create(
                content_type=PortalPublication.ContentType.IP_APPLICATION,
                object_id=application.id,
                defaults={
                    'is_public': application.status in {
                        IntellectualPropertyApplication.Status.AUTHORIZED,
                        IntellectualPropertyApplication.Status.ARCHIVED,
                    },
                    'is_featured': index == 0,
                    'member_consent': False,
                    'display_order': index,
                    'custom_title': application.title,
                    'custom_summary': application.intro,
                    'updated_by': teacher,
                },
            )

        governed_members = (
            self.leaders + self.primary_contributors + self.members[:10]
        )
        for index, member in enumerate(governed_members):
            has_consent = index < 12
            PortalPublication.objects.update_or_create(
                content_type=PortalPublication.ContentType.MEMBER,
                object_id=member.id,
                defaults={
                    'is_public': has_consent,
                    'is_featured': index < 4,
                    'member_consent': has_consent,
                    'display_order': index,
                    'custom_title': member.name,
                    'custom_summary': (
                        f'{DEMO_MARKER}{member.major} · {member.grade}，'
                        '参与团队项目协作与成果整理。'
                    ),
                    'updated_by': member if has_consent else teacher,
                },
            )

        self.stdout.write(self.style.SUCCESS(
            f'   门户治理创建完成：项目 {len(self.projects)} 项，'
            f'知识产权 {len(ip_applications)} 项，成员 {len(governed_members)} 人'
        ))

    # ------------------------------------------------------------------
    # 敏感资料
    # ------------------------------------------------------------------
    def create_sensitive_data(self):
        self.stdout.write('-> 创建敏感资料...')
        # 为 leader1, leader2, member1 各创建 1 条身份证敏感资料
        owners = [self.users['leader1'], self.users['leader2'], self.users['member1']]
        # 模拟身份证号（非真实号码，均为 110101 开头的演示号码）
        fake_ids = ['110101200001011234', '110101200106072345', '110101200311225678']

        self.sensitive_data_map = {}
        for index, (owner, fake_id) in enumerate(zip(owners, fake_ids)):
            sd = SensitiveData(
                data_type=SensitiveData.DataType.ID_CARD,
                title=f'{DEMO_MARKER}{owner.name}的身份证号码',
                display_name='身份证号码',
                key_version=1,
                project=self.projects[index],
                file_attachment=(
                    self.files_by_project[self.projects[0].id][2]
                    if index == 0
                    else None
                ),
                uploader=owner,
                is_encrypted=False,
                encrypted_content='',
            )
            sd.save()
            # 使用模型的加密方法加密模拟身份证号
            sd.encrypt_content(fake_id)
            self.sensitive_data_map[owner.email] = sd

        self.stdout.write(self.style.SUCCESS(
            f'   敏感资料创建完成：{len(owners)} 条（已加密）'
        ))

    # ------------------------------------------------------------------
    # 敏感资料访问申请
    # ------------------------------------------------------------------
    def create_sensitive_requests(self):
        self.stdout.write('-> 创建敏感资料访问申请...')
        now = timezone.now()
        approver = self.teachers[0]

        # 1. member5 申请下载 leader1 的附件 - 已通过，仍在有效期内
        sd_leader1 = self.sensitive_data_map[self.users['leader1'].email]
        SensitiveAccessRequest.objects.create(
            sensitive_data=sd_leader1,
            applicant=self.users['member5'],
            reason='办理比赛获奖奖金发放，需下载核验负责人证明附件。',
            usage_scenario='奖金发放附件核验',
            project=self.projects[0],
            expected_use_time=now + timedelta(hours=2),
            is_download=True,
            status=SensitiveAccessRequest.Status.APPROVED,
            approver=approver,
            approval_opinion='情况属实，同意在有效期内下载一次。',
            approved_at=now - timedelta(minutes=10),
            access_expires_at=now + timedelta(hours=2),
        )

        # 2. member8 只查看 leader2 的资料 - 已通过、不可下载
        sd_leader2 = self.sensitive_data_map[self.users['leader2'].email]
        SensitiveAccessRequest.objects.create(
            sensitive_data=sd_leader2,
            applicant=self.users['member8'],
            reason='赛前报名材料复核，仅需查看脱敏资料。',
            usage_scenario='报名材料线上复核',
            project=self.projects[1],
            expected_use_time=now + timedelta(hours=1),
            is_download=False,
            status=SensitiveAccessRequest.Status.APPROVED,
            approver=approver,
            approval_opinion='仅授权在线查看，不允许下载。',
            approved_at=now - timedelta(minutes=25),
            access_expires_at=now + timedelta(minutes=45),
            viewed_at=now - timedelta(minutes=5),
        )

        # 3. member10 的下载申请 - 待审批
        sd_member1 = self.sensitive_data_map[self.users['member1'].email]
        SensitiveAccessRequest.objects.create(
            sensitive_data=sd_member1,
            applicant=self.users['member10'],
            reason='协助办理报销需要核对经办人身份信息。',
            usage_scenario='财务报销身份核对',
            project=self.projects[2],
            expected_use_time=now + timedelta(days=1),
            is_download=True,
            status=SensitiveAccessRequest.Status.PENDING,
        )

        # 4. 历史查看授权已过期，用于演示审批与有效期过滤
        SensitiveAccessRequest.objects.create(
            sensitive_data=sd_leader2,
            applicant=self.users['member12'],
            reason='历史归档材料核验。',
            usage_scenario='结项归档复核',
            project=self.projects[1],
            expected_use_time=now - timedelta(days=2),
            is_download=False,
            status=SensitiveAccessRequest.Status.EXPIRED,
            approver=approver,
            approval_opinion='历史授权已超过有效期。',
            approved_at=now - timedelta(days=3),
            access_expires_at=now - timedelta(days=2),
            viewed_at=now - timedelta(days=2, hours=1),
        )

        self.stdout.write(self.style.SUCCESS('   敏感资料访问申请创建完成：4 条'))

    # ------------------------------------------------------------------
    # 公告
    # ------------------------------------------------------------------
    def create_announcements(self):
        self.stdout.write('-> 创建团队公告...')
        now = timezone.now()
        announcements = [
            {
                'title': '团队周例会与项目更新提醒',
                'content': (
                    '各项目负责人请在例会前完成项目进展更新；连续 11 天未更新'
                    '的项目将进入工作台提醒。'
                ),
                'category': Announcement.Category.PROJECT,
                'is_pinned': True,
                'is_public': False,
                'author': self.users['admin'],
            },
            {
                'title': '暑期项目路演安排',
                'content': (
                    '本周六下午进行阶段路演，请各项目准备一页进展摘要和'
                    '下一阶段计划。'
                ),
                'category': Announcement.Category.ACTIVITY,
                'is_pinned': False,
                'is_public': True,
                'author': self.users['teacher1'],
            },
            {
                'title': '报销票据 OCR 使用说明',
                'content': (
                    '上传清晰的票据图片后可自动识别日期、金额和票据号码，'
                    '提交前请人工核对识别结果。'
                ),
                'category': Announcement.Category.SYSTEM,
                'is_pinned': False,
                'is_public': False,
                'author': self.users['teacher2'],
            },
        ]
        for index, definition in enumerate(announcements):
            Announcement.objects.create(
                title=f'{DEMO_MARKER}{definition["title"]}',
                content=definition['content'],
                category=definition['category'],
                status=Announcement.Status.PUBLISHED,
                is_pinned=definition['is_pinned'],
                is_public=definition['is_public'],
                author=definition['author'],
                published_at=now - timedelta(days=index),
            )

        self.stdout.write(self.style.SUCCESS('   团队公告创建完成：3 条'))

    # ------------------------------------------------------------------
    # 通知
    # ------------------------------------------------------------------
    def create_notifications(self):
        self.stdout.write('-> 创建通知...')
        now = timezone.now()
        NT = Notification.NotificationType
        admin = self.users['admin']

        notifs = [
            {
                'recipient': self.users['member5'],
                'sender': admin,
                'type': NT.TASK,
                'title': '任务逾期提醒',
                'content': '您有任务已逾期，请尽快处理或申请延期。',
                'priority': Notification.Priority.HIGH,
                'ref_type': 'task', 'ref_id': 1,
            },
            {
                'recipient': self.users['leader1'],
                'sender': self.teachers[0],
                'type': NT.PROJECT,
                'title': '负责人更新提醒',
                'content': '您负责的项目已超过 11 天未更新进度，请及时填写进展。',
                'priority': Notification.Priority.NORMAL,
                'ref_type': 'project', 'ref_id': self.projects[2].id,
            },
            {
                'recipient': self.users['leader2'],
                'sender': admin,
                'type': NT.PROJECT,
                'title': '知识产权退回通知',
                'content': '您的发明专利申请被科研处退回，请尽快查看退回原因并修改。',
                'priority': Notification.Priority.HIGH,
                'ref_type': 'ip_application', 'ref_id': 2,
            },
            {
                'recipient': self.users['leader1'],
                'sender': self.users['member3'],
                'type': NT.PROJECT,
                'title': '贡献记录待审核',
                'content': '有新的贡献记录等待您审核，请及时处理。',
                'priority': Notification.Priority.NORMAL,
                'ref_type': 'contribution', 'ref_id': 1,
            },
            {
                'recipient': self.teachers[0],
                'sender': self.users['member10'],
                'type': NT.SYSTEM,
                'title': '敏感资料访问申请待审批',
                'content': '有新的敏感资料访问申请等待您审批。',
                'priority': Notification.Priority.HIGH,
                'ref_type': 'sensitive_request', 'ref_id': 2,
            },
            {
                'recipient': self.users['leader3'],
                'sender': admin,
                'type': NT.TASK,
                'title': '任务待审核',
                'content': '有任务已提交完成，等待您审核确认。',
                'priority': Notification.Priority.NORMAL,
                'ref_type': 'task', 'ref_id': 3,
            },
            {
                'recipient': self.users['member1'],
                'sender': admin,
                'type': NT.FINANCE,
                'title': '经费报销进度通知',
                'content': '您提交的报销申请已进入审核流程，请关注后续状态。',
                'priority': Notification.Priority.LOW,
                'ref_type': 'finance_expense', 'ref_id': 1,
            },
            {
                'recipient': self.users['leader1'],
                'sender': self.teachers[0],
                'type': NT.COMPETITION,
                'title': '比赛报名截止提醒',
                'content': '互联网+大赛报名即将截止，请尽快完成材料提交。',
                'priority': Notification.Priority.HIGH,
                'ref_type': 'competition', 'ref_id': 1,
            },
            {
                'recipient': self.users['member7'],
                'sender': admin,
                'type': NT.SYSTEM,
                'title': '灵活工时填写提醒',
                'content': '本周期灵活工作时间表尚未填写，请尽快提交。',
                'priority': Notification.Priority.NORMAL,
                'ref_type': 'work_schedule', 'ref_id': 1,
            },
            {
                'recipient': self.users['leader1'],
                'sender': self.users['member4'],
                'type': NT.PROJECT,
                'title': '排名异议通知',
                'content': '有成员对项目排名结果提出异议，请及时处理。',
                'priority': Notification.Priority.NORMAL,
                'ref_type': 'ranking_objection', 'ref_id': 1,
            },
        ]

        for n in notifs:
            Notification.objects.create(
                recipient=n['recipient'],
                sender=n['sender'],
                notification_type=n['type'],
                channel=Notification.Channel.INAPP,
                title=f'{DEMO_MARKER}{n["title"]}',
                content=n['content'],
                priority=n['priority'],
                related_object_type=n['ref_type'],
                related_object_id=n['ref_id'],
                is_read=random.choice([True, False, False]),
            )

        self.stdout.write(self.style.SUCCESS(f'   通知创建完成：{len(notifs)} 条'))

    # ------------------------------------------------------------------
    # 操作日志
    # ------------------------------------------------------------------
    def create_operation_logs(self):
        self.stdout.write('-> 创建操作日志...')
        OT = OperationLog.OperationType

        log_specs = [
            (self.users['admin'], OT.LOGIN, 'auth', 'User', '登录系统', 'POST', '/api/auth/login/', 200),
            (self.users['leader1'], OT.CREATE, 'projects', 'Project', '创建项目 智能校园导览系统', 'POST', '/api/projects/', 201),
            (self.users['leader1'], OT.UPDATE, 'projects', 'Project', '更新项目进度', 'PATCH', '/api/projects/1/', 200),
            (self.users['member3'], OT.CREATE, 'tasks', 'Task', '创建任务 需求调研与整理', 'POST', '/api/tasks/', 201),
            (self.users['member5'], OT.UPDATE, 'tasks', 'Task', '更新任务状态为进行中', 'PATCH', '/api/tasks/3/', 200),
            (self.users['leader2'], OT.APPROVE, 'contributions', 'Contribution', '审核通过贡献记录', 'POST', '/api/contributions/1/review/', 200),
            (self.users['leader1'], OT.REJECT, 'contributions', 'Contribution', '驳回贡献记录', 'POST', '/api/contributions/2/review/', 200),
            (self.users['member10'], OT.CREATE, 'sensitive', 'SensitiveAccessRequest', '提交敏感资料访问申请', 'POST', '/api/sensitive/requests/', 201),
            (self.teachers[0], OT.APPROVE, 'sensitive', 'SensitiveAccessRequest', '审批通过敏感资料访问申请', 'POST', '/api/sensitive/requests/1/approve/', 200),
            (self.users['member5'], OT.VIEW_SENSITIVE, 'sensitive', 'SensitiveData', '查看敏感资料明文', 'GET', '/api/sensitive/1/view/', 200),
            (self.users['admin'], OT.UPLOAD, 'files', 'FileAsset', '上传项目文件', 'POST', '/api/files/', 201),
            (self.users['member12'], OT.DOWNLOAD, 'files', 'FileAsset', '下载公开文件', 'GET', '/api/files/1/download/', 200),
            (self.users['leader3'], OT.CREATE, 'ip', 'IPApplication', '创建知识产权申请', 'POST', '/api/ip/applications/', 201),
            (self.users['teacher1'], OT.REVIEW, 'ip', 'IPApplication', '确认知识产权申请', 'POST', '/api/ip/applications/1/confirm/', 200),
            (self.users['admin'], OT.EXPORT, 'exports', 'Project', '导出项目汇总报告', 'GET', '/api/exports/project-report/', 200),
            (self.users['leader1'], OT.CREATE, 'finance', 'FinanceExpense', '登记经费支出', 'POST', '/api/finance/expenses/', 201),
            (self.users['member8'], OT.DELETE, 'tasks', 'Task', '删除作废任务', 'DELETE', '/api/tasks/9/', 204),
            (self.users['admin'], OT.IMPORT, 'imports', 'Member', '批量导入成员', 'POST', '/api/imports/members/', 200),
            (self.users['member15'], OT.UPDATE, 'members', 'FlexibleWorkSchedule', '填写灵活工时', 'POST', '/api/members/work-schedules/', 201),
            (self.users['leader2'], OT.OTHER, 'rankings', 'MemberRanking', '提交成员排序草案', 'POST', '/api/contributions/rankings/', 201),
        ]

        for spec in log_specs:
            operator, op_type, module, obj_type, desc, method, path, status = spec
            OperationLog.objects.create(
                operator=operator,
                operation_type=op_type,
                module=module,
                object_type=obj_type,
                object_id=str(random.randint(1, 50)),
                description=f'{DEMO_MARKER}{desc}',
                request_method=method,
                request_path=path,
                request_ip='127.0.0.1',
                user_agent='Mozilla/5.0 (Windows NT 10.0; Win64; x64) DemoBrowser/1.0',
                request_data={'demo': True, 'source': 'seed_demo_data'},
                response_status=status,
                is_success=status < 400,
            )

        self.stdout.write(self.style.SUCCESS(f'   操作日志创建完成：{len(log_specs)} 条'))

    # ------------------------------------------------------------------
    # 清除演示数据
    # ------------------------------------------------------------------
    def clean_demo_data(self):
        self.stdout.write(self.style.WARNING('-> 精准清除本命令的演示数据...'))
        demo_projects = Project.all_objects.filter(
            code__startswith=DEMO_PROJECT_PREFIX
        )
        legacy_projects = Project.all_objects.filter(
            code__startswith='PROJ-',
            leader__email__endswith='@demo.com',
        )
        legacy_competition_projects = Project.all_objects.filter(
            code__startswith='DEMO-2026-',
        )
        project_ids = list(
            demo_projects.values_list('id', flat=True)
        ) + list(
            legacy_projects.values_list('id', flat=True)
        ) + list(
            legacy_competition_projects.values_list('id', flat=True)
        )
        project_ids = list(dict.fromkeys(project_ids))
        competition_queryset = Competition.objects.filter(
            project_id__in=project_ids
        )
        competition_event_ids = list(dict.fromkeys([
            *CompetitionEvent.objects.filter(
                organization__code=DEMO_TEAM_CODE
            ).values_list('id', flat=True),
            *list(
                competition_queryset.exclude(event_id=None).values_list(
                    'event_id',
                    flat=True,
                ).distinct()
            ),
        ]))
        ip_applications = IntellectualPropertyApplication.objects.filter(
            application_code__startswith=DEMO_IP_PREFIX
        )
        if project_ids:
            ip_applications = (
                ip_applications
                | IntellectualPropertyApplication.objects.filter(
                    related_project_id__in=project_ids
                )
            )
        ip_ids = list(
            ip_applications.values_list('id', flat=True).distinct()
        )

        total = 0

        def delete_queryset(queryset):
            nonlocal total
            count, _ = queryset.delete()
            total += count

        # 删除数据库记录前先清除本命令写入的物理文件。
        receipt_queryset = FinanceReceipt.objects.filter(
            expense__project_id__in=project_ids
        )
        for receipt in receipt_queryset.exclude(file='').iterator():
            receipt.file.delete(save=False)

        version_queryset = FileVersion.objects.filter(
            file_asset__project_id__in=project_ids
        )
        for version in version_queryset.exclude(file='').iterator():
            version.file.delete(save=False)
        file_queryset = FileAsset.objects.filter(project_id__in=project_ids)
        for asset in file_queryset.exclude(file='').iterator():
            asset.file.delete(save=False)

        demo_report_queryset = CustomReport.objects.filter(
            Q(name__startswith=DEMO_MARKER)
            | Q(
                name__in=LEGACY_COMPETITION_REPORT_NAMES,
                created_by__email__endswith='@demo.com',
            )
        )
        report_execution_queryset = ScheduledReportExecution.objects.filter(
            schedule__report__in=demo_report_queryset
        )
        for execution in report_execution_queryset.exclude(file='').iterator():
            execution.file.delete(save=False)

        import_queryset = ImportTask.objects.filter(
            file_path__contains=str(
                Path('imports') / DEMO_IMPORT_DIRNAME
            )
        )
        safe_import_root = (
            Path(settings.MEDIA_ROOT)
            / 'imports'
            / DEMO_IMPORT_DIRNAME
        ).resolve()
        for file_path in import_queryset.values_list('file_path', flat=True):
            try:
                candidate = Path(file_path).resolve()
                if candidate.is_relative_to(safe_import_root):
                    candidate.unlink(missing_ok=True)
            except (OSError, RuntimeError):
                pass
        delete_queryset(import_queryset)

        if project_ids:
            delete_queryset(PortalPublication.objects.filter(
                content_type=PortalPublication.ContentType.PROJECT,
                object_id__in=project_ids,
            ))
        if ip_ids:
            delete_queryset(PortalPublication.objects.filter(
                content_type=PortalPublication.ContentType.IP_APPLICATION,
                object_id__in=ip_ids,
            ))
        delete_queryset(PortalPublication.objects.filter(
            content_type=PortalPublication.ContentType.MEMBER,
            custom_summary__startswith=DEMO_MARKER,
        ))
        delete_queryset(PortalSettings.objects.filter(
            singleton_key='default',
            team_name='数智创新实践团队',
            updated_by__email__endswith='@demo.com',
        ))

        delete_queryset(OperationLog.objects.filter(
            description__startswith=DEMO_MARKER
        ))
        delete_queryset(OperationLog.objects.filter(
            request_data__source='seed_demo_data'
        ))
        delete_queryset(Notification.objects.filter(
            title__startswith=DEMO_MARKER
        ))
        delete_queryset(Announcement.objects.filter(
            title__startswith=DEMO_MARKER
        ))
        delete_queryset(demo_report_queryset)
        delete_queryset(UserLifecycleEvent.objects.filter(
            reason__startswith=DEMO_MARKER
        ))
        # 演示账号会被其他旧演示脚本复用；完整重建时清除其技能关系，
        # 避免旧脚本的随机技能被并入本命令的确定性演示包。
        # SkillTag 仍作为共享字典保留，真实账号的技能关系不受影响。
        delete_queryset(MemberSkill.objects.filter(
            user__email__in=DEMO_ACCOUNT_EMAILS
        ))
        delete_queryset(FlexibleWorkSchedule.objects.filter(
            notes__startswith=DEMO_MARKER
        ))

        if project_ids:
            delete_queryset(SensitiveAccessRequest.objects.filter(
                project_id__in=project_ids
            ))
        delete_queryset(SensitiveData.objects.filter(
            title__startswith=DEMO_MARKER
        ))

        if ip_ids:
            delete_queryset(IPObjection.objects.filter(
                application_id__in=ip_ids
            ))
            delete_queryset(IPReturnRecord.objects.filter(
                application_id__in=ip_ids
            ))
            delete_queryset(IPApplicationContributor.objects.filter(
                application_id__in=ip_ids
            ))
            delete_queryset(
                IntellectualPropertyApplication.objects.filter(id__in=ip_ids)
            )

        if project_ids:
            delete_queryset(RankingObjection.objects.filter(
                ranking__project_id__in=project_ids
            ))
            delete_queryset(MemberRanking.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(Contribution.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(receipt_queryset)
            delete_queryset(FinanceExpense.all_objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(FinanceIncome.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(FinanceBudget.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(Task.all_objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(version_queryset)
            delete_queryset(file_queryset)
            delete_queryset(competition_queryset)
            delete_queryset(CompetitionEvent.objects.filter(
                id__in=competition_event_ids,
                entries__isnull=True,
            ))
            delete_queryset(ProjectMembershipEvent.objects.filter(
                membership__project_id__in=project_ids
            ))
            delete_queryset(ProjectStageLog.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(ProjectMember.objects.filter(
                project_id__in=project_ids
            ))
            delete_queryset(Project.all_objects.filter(id__in=project_ids))

        # 旧版固定小组必须在根团队前删除；当前小团队版不再创建固定部门。
        delete_queryset(Team.objects.filter(code__in=LEGACY_DEMO_SQUAD_CODES))
        delete_queryset(Team.objects.filter(code=DEMO_TEAM_CODE))

        # 精确清理旧版专属账号，避免缩减成员规模后残留 member36~52、
        # 旧审批员或旧比赛负责人。
        delete_queryset(User.objects.filter(
            email__in=(
                *LEGACY_COMPETITION_ACCOUNT_EMAILS,
                *LEGACY_DEMO_ACCOUNT_EMAILS,
            )
        ))

        # 当前账号和共享技能标签可被其他演示场景复用，不做删除；创建阶段使用
        # update_or_create 保证重复运行稳定。所有真实账号和非演示业务记录均不触碰。
        self.stdout.write(self.style.SUCCESS(
            f'   已精准清除 {total} 条团队演示记录；真实数据与共享演示账号均保留'
        ))

    # ------------------------------------------------------------------
    # 账号汇总
    # ------------------------------------------------------------------
    def print_account_summary(self):
        self.stdout.write('账号清单：')
        self.stdout.write('  系统管理员: admin@demo.com / admin123456')
        teacher1_label = (
            '操作老师'
            if getattr(self, 'demo_teacher_is_operator', True)
            else '查看老师'
        )
        self.stdout.write(
            f'  {teacher1_label}:  teacher1@demo.com / teacher123456'
        )
        self.stdout.write(
            '  查看老师:  teacher2~teacher4@demo.com / teacher123456'
        )
        self.stdout.write(
            '  五位负责人: leader1~leader5@demo.com / leader123456'
        )
        self.stdout.write(
            '  平行主贡献者: contributor1~contributor2@demo.com / member123456'
        )
        self.stdout.write(
            '  普通成员:  member1~35@demo.com / member123456'
        )
