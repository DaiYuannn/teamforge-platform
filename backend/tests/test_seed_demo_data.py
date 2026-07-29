from collections import Counter
from io import BytesIO
from io import StringIO
from pathlib import Path

import pytest
from docx import Document
from django.core.files.base import ContentFile
from django.core.files.uploadedfile import SimpleUploadedFile
from django.core.management import call_command
from django.utils import timezone
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from apps.common.team_models import Team, TeamMember, TeamMembershipEvent
from apps.competitions.models import (
    Competition,
    CompetitionAward,
    CompetitionEvent,
    CompetitionParticipant,
)
from apps.contributions.models import Contribution
from apps.dashboard.portal_models import PortalPublication, PortalSettings
from apps.exports.custom_report_models import CustomReport
from apps.exports.scheduled_report_models import (
    ScheduledReport,
    ScheduledReportExecution,
)
from apps.files.models import FileAsset, FileVersion
from apps.finance.models import FinanceExpense, FinanceIncome, FinanceReceipt
from apps.finance.ocr_service import parse_receipt_text, validate_image
from apps.imports.models import ImportTask
from apps.intellectual_property.models import (
    IPApplicationContributor,
    IPMaterialVersion,
    IntellectualPropertyApplication,
)
from apps.members.models import MemberSkill, SkillTag
from apps.notifications.models import Announcement
from apps.projects.models import (
    Project,
    ProjectMember,
    ProjectMembershipEvent,
    ProjectStageLog,
)
from apps.sensitive.models import SensitiveAccessRequest, SensitiveData
from apps.tasks.models import Task
from apps.users.management.commands.seed_demo_data import (
    CORE_COMPETITION_SERIES,
    DEMO_ACCOUNT_EMAILS,
    DEMO_IMPORT_DIRNAME,
    DEMO_IP_PREFIX,
    DEMO_MARKER,
    DEMO_PROJECT_PREFIX,
    DEMO_SQUAD_CODES,
    DEMO_TEAM_CODE,
    LEGACY_COMPETITION_ACCOUNT_EMAILS,
    LEGACY_DEMO_ACCOUNT_EMAILS,
    LEGACY_COMPETITION_REPORT_NAMES,
    SELECTIVE_COMPETITION_SERIES,
)
from apps.users.models import User, UserLifecycleEvent, UserPreference


@pytest.mark.django_db(transaction=True)
def test_seed_demo_data_scale_lifecycle_assets_and_safe_clean(settings, tmp_path):
    settings.MEDIA_ROOT = tmp_path / 'media'

    real_user = User.objects.create_user(
        email='owner@real.example',
        username='real-owner',
        password='safe-password',
        name='真实团队负责人',
        global_role=User.GlobalRole.TEACHER,
        is_student=False,
    )
    real_team = Team.objects.create(
        name='真实团队不可删除',
        code='REAL-TEAM-KEEP',
        owner=real_user,
    )
    real_team_member = TeamMember.objects.create(
        team=real_team,
        user=real_user,
        role=TeamMember.Role.OWNER,
    )
    real_team_event = TeamMembershipEvent.objects.create(
        membership=real_team_member,
        event_type='joined',
        to_role=TeamMember.Role.OWNER,
        to_status=TeamMember.Status.ACTIVE,
        reason='真实团队成员加入记录',
        operator=real_user,
    )
    real_project = Project.objects.create(
        name='真实项目不可删除',
        code='REAL-KEEP-001',
        leader=real_user,
        current_stage=Project.Stage.DEV_EXPERIMENT,
        start_date=timezone.localdate(),
        status=Project.Status.ACTIVE,
    )
    legacy_demo_project = Project.objects.create(
        name='旧版演示项目应清理',
        code='DEMO-2026-999',
        leader=real_user,
        current_stage=Project.Stage.DEV_EXPERIMENT,
        start_date=timezone.localdate(),
        status=Project.Status.ACTIVE,
    )
    legacy_demo_ip = IntellectualPropertyApplication.objects.create(
        title='旧版演示知识产权应清理',
        application_code='IP-DEMO-2026-999',
        related_project=legacy_demo_project,
        main_writer=real_user,
        created_by=real_user,
    )
    ProjectMember.objects.create(
        project=real_project,
        user=real_user,
        role_in_project=ProjectMember.RoleInProject.LEADER,
    )
    real_task = Task.objects.create(
        project=real_project,
        title='真实任务不可删除',
        assignee=real_user,
        creator=real_user,
    )
    real_income = FinanceIncome.objects.create(
        project=real_project,
        title='真实项目收入',
        amount='1000.00',
        income_type=FinanceIncome.IncomeType.GRANT,
        income_date=timezone.localdate(),
        source='真实来源',
        reference_number='REAL-INCOME-001',
        recorded_by=real_user,
    )
    real_file = FileAsset(
        project=real_project,
        name='真实文件不可删除.txt',
        level=FileAsset.Level.INTERNAL,
        size=12,
        content_type='text/plain',
        uploader=real_user,
    )
    real_file.file.save(
        'real/keep.txt',
        ContentFile(b'keep forever'),
        save=True,
    )
    real_import = ImportTask.objects.create(
        module=ImportTask.Module.PROJECTS,
        file_path=str(tmp_path / 'real-import.xlsx'),
        status=ImportTask.Status.PREVIEWED,
        created_by=real_user,
    )
    real_skill = SkillTag.objects.create(name='真实团队专属技能')
    portal_settings = PortalSettings.objects.create(
        singleton_key='default',
        team_name='真实团队门户',
        updated_by=real_user,
    )
    real_publication = PortalPublication.objects.create(
        content_type=PortalPublication.ContentType.PROJECT,
        object_id=real_project.id,
        is_public=True,
        updated_by=real_user,
    )
    real_lifecycle = UserLifecycleEvent.objects.create(
        user=real_user,
        event_type=UserLifecycleEvent.EventType.CREATED,
        to_status=User.MembershipStatus.ACTIVE,
        reason='真实成员加入记录',
        operator=real_user,
    )
    real_report = CustomReport.objects.create(
        name='真实团队月报不可删除',
        description='真实团队报表',
        report_type=CustomReport.ReportType.SUMMARY,
        config={'data_source': 'project', 'group_by': 'status'},
        created_by=real_user,
        is_scheduled=True,
    )
    real_schedule = ScheduledReport.objects.create(
        report=real_report,
        created_by=real_user,
        frequency=ScheduledReport.Frequency.MONTHLY,
        execution_time=timezone.localtime().time().replace(microsecond=0),
        file_format=ScheduledReport.FileFormat.XLSX,
    )
    real_execution = ScheduledReportExecution(
        schedule=real_schedule,
        status=ScheduledReport.RunStatus.SUCCESS,
        file_name='real-team-report.xlsx',
        file_format=ScheduledReport.FileFormat.XLSX,
        file_size=16,
        finished_at=timezone.now(),
        generated_by=real_user,
    )
    real_execution.file.save(
        'real/keep-report.xlsx',
        ContentFile(b'real report data'),
        save=True,
    )
    real_announcement = Announcement.objects.create(
        title='真实团队公告不可删除',
        content='真实公告内容',
        category=Announcement.Category.SYSTEM,
        status=Announcement.Status.PUBLISHED,
        author=real_user,
        published_at=timezone.now(),
    )
    for legacy_email in LEGACY_DEMO_ACCOUNT_EMAILS:
        User.objects.create_user(
            email=legacy_email,
            username=legacy_email.split('@', 1)[0],
            password='legacy-demo-password',
            name='旧版小团队演示账号',
        )

    output = StringIO()
    call_command('seed_demo_data', clean=True, force=True, stdout=output)

    assert not Project.all_objects.filter(pk=legacy_demo_project.pk).exists()
    assert not IntellectualPropertyApplication.objects.filter(
        pk=legacy_demo_ip.pk
    ).exists()
    assert not User.objects.filter(
        email__in=LEGACY_DEMO_ACCOUNT_EMAILS
    ).exists()

    owned_emails = list(DEMO_ACCOUNT_EMAILS)
    demo_projects = Project.all_objects.filter(
        code__startswith=DEMO_PROJECT_PREFIX
    )
    demo_project_ids = list(demo_projects.values_list('id', flat=True))
    demo_tasks = Task.all_objects.filter(project_id__in=demo_project_ids)
    demo_files = FileAsset.objects.filter(project_id__in=demo_project_ids)

    assert User.objects.filter(email__in=owned_emails).count() == 47
    assert User.objects.filter(
        email__in=owned_emails,
        global_role=User.GlobalRole.MEMBER,
    ).count() == 46
    assert User.objects.filter(
        email__in=owned_emails,
        global_role=User.GlobalRole.TEACHER,
    ).count() == 0
    assert User.objects.filter(
        global_role=User.GlobalRole.TEACHER,
        is_active=True,
    ).count() == 1
    assert User.objects.get(email='teacher1@demo.com').global_role == (
        User.GlobalRole.MEMBER
    )
    assert User.objects.get(email='teacher2@demo.com').global_role == (
        User.GlobalRole.MEMBER
    )
    assert set(
        User.objects.filter(email__in=owned_emails).values_list(
            'membership_status', flat=True
        )
    ) == {
        User.MembershipStatus.ACTIVE,
        User.MembershipStatus.ON_LEAVE,
        User.MembershipStatus.EXITED,
        User.MembershipStatus.EXTERNAL,
    }
    exited_users = User.objects.filter(
        email__in=owned_emails,
        membership_status=User.MembershipStatus.EXITED,
    )
    assert exited_users.count() == 2
    assert exited_users.filter(is_active=False).count() == 2
    assert User.objects.filter(
        email__in=owned_emails,
        membership_status__in=[
            User.MembershipStatus.ACTIVE,
            User.MembershipStatus.ON_LEAVE,
            User.MembershipStatus.EXTERNAL,
        ],
        is_active=True,
    ).count() == 45

    demo_team = Team.objects.get(code=DEMO_TEAM_CODE)
    demo_team_members = TeamMember.objects.filter(team=demo_team)
    assert demo_team_members.count() == 46
    assert Counter(demo_team_members.values_list('role', flat=True)) == {
        TeamMember.Role.OWNER: 1,
        TeamMember.Role.CO_LEAD: 6,
        TeamMember.Role.TEACHER: 4,
        TeamMember.Role.MEMBER: 34,
        TeamMember.Role.EXTERNAL: 1,
    }
    assert Counter(demo_team_members.values_list('status', flat=True)) == {
        TeamMember.Status.ACTIVE: 41,
        TeamMember.Status.ON_LEAVE: 3,
        TeamMember.Status.EXITED: 2,
    }
    demo_team_events = TeamMembershipEvent.objects.filter(
        membership__team=demo_team
    )
    assert demo_team_events.count() == 53
    assert Counter(demo_team_events.values_list('event_type', flat=True)) == {
        'joined': 46,
        'status_changed': 3,
        'exited': 2,
        'handover': 2,
    }
    assert demo_team_events.filter(reason__startswith=DEMO_MARKER).count() == 53
    demo_squads = Team.objects.filter(parent=demo_team).order_by('code')
    assert demo_squads.count() == 0
    assert DEMO_SQUAD_CODES == ()
    assert Team.objects.filter(pk=real_team.pk).exists()

    preferences = UserPreference.objects.filter(user__email__in=owned_emails)
    assert preferences.count() == 47
    assert set(preferences.values_list('items_per_page', flat=True)) <= {10, 20, 50}
    assert set(preferences.values_list('theme_mode', flat=True)) == {'system'}
    assert set(preferences.values_list('schedule_start', flat=True)) == {'19:00'}
    assert set(preferences.values_list('schedule_end', flat=True)) == {'07:00'}
    expected_preference_keys = {
        'system', 'task', 'project', 'competition', 'finance',
        'contribution', 'schedule', 'approval', 'report',
    }
    for preference in preferences:
        assert len(preference.sidebar_order) == 5
        assert preference.favorite_routes
        assert preference.saved_filters
        assert set(
            preference.notification_preferences['categories']
        ) == expected_preference_keys
        assert set(
            preference.notification_preferences['channels']
        ) == {'in_app', 'email'}
        assert set(
            preference.notification_preferences['quiet_hours']
        ) == {'enabled', 'start', 'end'}
        assert preference.notification_preferences['digest'] in {
            'instant', 'daily', 'weekly',
        }

    admin_preference = UserPreference.objects.get(user__email='admin@demo.com')
    member_preference = UserPreference.objects.get(user__email='member1@demo.com')
    external_preference = UserPreference.objects.get(user__email='member35@demo.com')
    allowed_dashboard_cards = {'signals', 'priority', 'delivery', 'business'}
    assert set(admin_preference.dashboard_layout['cards']) == allowed_dashboard_cards
    assert set(member_preference.dashboard_layout['cards']) == allowed_dashboard_cards
    assert set(external_preference.dashboard_layout['cards']) == allowed_dashboard_cards
    assert len({
        tuple(admin_preference.dashboard_layout['cards']),
        tuple(member_preference.dashboard_layout['cards']),
        tuple(external_preference.dashboard_layout['cards']),
    }) == 3
    assert admin_preference.dashboard_layout['profile'] == 'admin'
    assert admin_preference.notification_preferences['digest'] == 'instant'
    assert member_preference.dashboard_layout['profile'] == 'member'
    assert external_preference.dashboard_layout['profile'] == 'external'
    assert external_preference.notification_preferences['digest'] == 'weekly'
    assert external_preference.favorite_routes == ['/tasks', '/files']

    assert demo_projects.count() == 7
    for project in demo_projects.prefetch_related('teams'):
        assert list(project.teams.all()) == [demo_team]
        assert project.visibility == Project.Visibility.PROJECT
        assert ProjectMember.objects.filter(
            project=project,
            user=project.leader,
            role_in_project=ProjectMember.RoleInProject.LEADER,
        ).exists()
    project_years = set(
        demo_projects.values_list('start_date__year', flat=True)
    )
    assert project_years == {2021, 2024}
    assert demo_projects.filter(status=Project.Status.CLOSED).count() == 2
    assert demo_projects.filter(status=Project.Status.ACTIVE).count() == 5
    for year in range(2021, 2024):
        active_that_year = [
            project
            for project in demo_projects
            if project.start_date.year <= year
            and (
                project.actual_end_date is None
                or project.actual_end_date.year >= year
            )
        ]
        assert len(active_that_year) <= 3
    active_in_2024 = [
        project
        for project in demo_projects
        if project.start_date.year <= 2024
        and (
            project.actual_end_date is None
            or project.actual_end_date.year >= 2024
        )
    ]
    assert len(active_in_2024) == 5
    assert demo_tasks.count() == 35
    assert not demo_tasks.filter(
        status__in=[Task.Status.PENDING_REVIEW, Task.Status.DONE],
        completion_note='',
    ).exists()
    assert not demo_tasks.filter(
        status__in=[Task.Status.PENDING_REVIEW, Task.Status.DONE],
        reviewer__isnull=True,
    ).exists()
    assert not demo_tasks.filter(
        status__in=[Task.Status.OVERDUE, Task.Status.NEED_HELP],
        delay_reason='',
    ).exists()
    assert ProjectMembershipEvent.objects.filter(
        membership__project_id__in=demo_project_ids
    ).count() == 71
    assert ProjectStageLog.objects.filter(
        project_id__in=demo_project_ids
    ).count() == 54
    assert UserLifecycleEvent.objects.filter(
        reason__startswith=DEMO_MARKER
    ).count() == 48
    assert {
        ProjectMember.Status.ACTIVE,
        ProjectMember.Status.ON_LEAVE,
        ProjectMember.Status.EXITED,
    }.issubset(
        set(
            ProjectMember.objects.filter(
                project_id__in=demo_project_ids
            ).values_list('status', flat=True)
        )
    )

    competitions = Competition.objects.filter(
        project_id__in=demo_project_ids
    ).select_related('event', 'project')
    competition_events = CompetitionEvent.objects.filter(
        organization=demo_team
    )
    assert competition_events.count() == 33
    assert competitions.count() == 109
    assert not competitions.filter(event__isnull=True).exists()
    assert len({
        (competition.event_id, competition.project_id)
        for competition in competitions
    }) == competitions.count()

    current_edition = str(timezone.localdate().year)
    for series in CORE_COMPETITION_SERIES:
        latest_entries = competitions.filter(
            event__name=series['name'],
            event__edition=current_edition,
        )
        assert latest_entries.count() == 5
        assert latest_entries.values('project_id').distinct().count() == 5
    for series in SELECTIVE_COMPETITION_SERIES:
        for event in competition_events.filter(name=series['name']):
            assert 1 <= event.entries.count() <= 3

    root_user_ids = set(
        demo_team_members.values_list('user_id', flat=True)
    )
    for competition in competitions:
        participants = list(
            competition.participants.select_related('user')
        )
        assert len(participants) == 6
        assert sum(
            participant.role == CompetitionParticipant.Role.LEADER
            for participant in participants
        ) == 1
        assert {
            participant.participation_status
            for participant in participants
        } == {
            CompetitionParticipant.ParticipationStatus.CONFIRMED
        }
        assert {
            participant.user_id for participant in participants
        }.issubset(root_user_ids)

    for award in CompetitionAward.objects.filter(
        competition__in=competitions
    ).prefetch_related('recipients', 'competition__participants'):
        confirmed_ids = set(
            award.competition.participants.filter(
                participation_status=(
                    CompetitionParticipant.ParticipationStatus.CONFIRMED
                )
            ).values_list('user_id', flat=True)
        )
        assert set(
            award.recipients.values_list('id', flat=True)
        ).issubset(confirmed_ids)
        assert award.competition.is_awarded
        assert award.competition.award_level == award.award_level

    retirement_awards = competitions.filter(
        award_level__in=['国赛金奖', '国赛银奖']
    )
    assert retirement_awards.count() == 2
    for retired_entry in retirement_awards:
        assert not competitions.filter(
            project=retired_entry.project,
            event__name=retired_entry.event.name,
            comp_type=retired_entry.comp_type,
            event__edition__gt=retired_entry.event.edition,
        ).exists()

    assert demo_files.count() == 247
    assert Counter(
        Path(name).suffix.lower()
        for name in demo_files.values_list('file', flat=True)
    ) == {
        '.pdf': 8,
        '.docx': 116,
        '.xlsx': 7,
        '.pptx': 116,
    }
    assert FileVersion.objects.filter(
        file_asset__project_id__in=demo_project_ids
    ).count() == 116
    assert Task.attachment_files.through.objects.filter(
        task__project_id__in=demo_project_ids
    ).count() == 28
    for asset in demo_files.select_related('project', 'project__leader'):
        content = Path(asset.file.path).read_bytes()
        assert asset.size == len(content)
        suffix = Path(asset.file.name).suffix.lower()
        if suffix == '.pdf':
            assert content.startswith(b'%PDF')
            parsed_text = '\n'.join(
                page.extract_text() or ''
                for page in PdfReader(BytesIO(content)).pages
            )
        elif suffix == '.docx':
            assert content.startswith(b'PK')
            parsed_text = '\n'.join(
                paragraph.text
                for paragraph in Document(BytesIO(content)).paragraphs
            )
        elif suffix == '.xlsx':
            assert content.startswith(b'PK')
            workbook = load_workbook(BytesIO(content), data_only=True)
            parsed_text = '\n'.join(
                str(cell.value or '')
                for row in workbook.active.iter_rows()
                for cell in row
            )
        else:
            assert content.startswith(b'PK')
            presentation = Presentation(BytesIO(content))
            assert len(presentation.slides) >= 1
            parsed_text = '\n'.join(
                shape.text
                for slide in presentation.slides
                for shape in slide.shapes
                if hasattr(shape, 'text')
            )
            assert parsed_text.strip()

        assert asset.project.code in parsed_text
        assert asset.project.name in parsed_text
        assert asset.project.leader.name in parsed_text

    competition_files = demo_files.filter(
        competition_entry__isnull=False
    ).select_related('competition_entry', 'competition_entry__event')
    assert competition_files.count() == 218
    competition_hashes = list(
        competition_files.values_list('file_hash', flat=True)
    )
    assert all(competition_hashes)
    assert len(set(competition_hashes)) == len(competition_hashes)
    for competition in competitions:
        entry_files = competition_files.filter(
            competition_entry=competition
        )
        assert entry_files.count() == 2
        assert {
            Path(name).suffix.lower()
            for name in entry_files.values_list('file', flat=True)
        } == {'.docx', '.pptx'}
        assert not entry_files.exclude(
            project_id=competition.project_id
        ).exists()
        plan_asset = entry_files.get(file__endswith='.docx')
        assert plan_asset.version == 2
        assert competition.event.edition in plan_asset.name
        assert competition.event.name in plan_asset.name
        assert plan_asset.created_at.year == int(competition.event.edition)
        historical_plan = FileVersion.objects.get(
            file_asset=plan_asset,
            version=1,
        )
        assert historical_plan.created_at.year == int(
            competition.event.edition
        )
        assert (
            Path(historical_plan.file.path).read_bytes()
            != Path(plan_asset.file.path).read_bytes()
        )

    demo_incomes = FinanceIncome.objects.filter(
        project_id__in=demo_project_ids
    ).select_related('project', 'competition_entry__event')
    bonus_incomes = demo_incomes.filter(
        income_type=FinanceIncome.IncomeType.BONUS
    )
    grant_incomes = demo_incomes.filter(
        income_type=FinanceIncome.IncomeType.GRANT
    )
    assert bonus_incomes.count() == demo_projects.count()
    assert grant_incomes.count() == demo_projects.count()
    assert not bonus_incomes.filter(competition_entry__isnull=True).exists()
    assert not grant_incomes.filter(competition_entry__isnull=False).exists()
    for bonus_income in bonus_incomes:
        assert (
            bonus_income.competition_entry.project_id
            == bonus_income.project_id
        )
        assert bonus_income.competition_entry.is_awarded

    receipts = FinanceReceipt.objects.filter(
        expense__project_id__in=demo_project_ids
    )
    demo_expenses = FinanceExpense.all_objects.filter(
        project_id__in=demo_project_ids
    )
    assert 14 <= demo_expenses.count() <= 21
    competition_expenses = demo_expenses.filter(
        competition_entry__isnull=False
    ).select_related('competition_entry__event')
    for expense in competition_expenses:
        assert expense.project_id == expense.competition_entry.project_id

    for project in demo_projects:
        travel_expense = competition_expenses.get(
            project=project,
            title='比赛现场往返交通费',
            category=FinanceExpense.Category.TRAVEL,
        )
        registration_expense = competition_expenses.get(
            project=project,
            title='比赛报名费',
            category=FinanceExpense.Category.COMPETITION_FEE,
        )
        assert (
            travel_expense.competition_entry_id
            != registration_expense.competition_entry_id
        )
        assert (
            travel_expense.competition_entry.event_id
            != registration_expense.competition_entry.event_id
        )
        if project.status == Project.Status.ACTIVE:
            assert {
                travel_expense.competition_entry.event.edition,
                registration_expense.competition_entry.event.edition,
            } == {current_edition}

    assert demo_expenses.count() <= receipts.count() <= demo_expenses.count() * 2
    expected_receipt_count = receipts.count()
    assert {
        Path(name).suffix.lower()
        for name in receipts.values_list('file', flat=True)
    } == {'.png'}
    first_receipt_content = None
    for receipt in receipts:
        content = Path(receipt.file.path).read_bytes()
        with Image.open(BytesIO(content)) as image:
            image.verify()
        first_receipt_content = first_receipt_content or content

    uploaded_receipt = SimpleUploadedFile(
        'demo-receipt.png',
        first_receipt_content,
        content_type='image/png',
    )
    validate_image(uploaded_receipt)
    parsed_receipt = parse_receipt_text(
        'DEMO SUPPLY STORE\nDATE: 2026-07-20\n'
        'TOTAL: 128.50\nNO: DEMO-000001',
        confidence=0.92,
    )
    assert parsed_receipt['vendor'] == 'DEMO SUPPLY STORE'
    assert parsed_receipt['expense_date'] == '2026-07-20'
    assert parsed_receipt['amount'] == '128.50'
    assert parsed_receipt['invoice_number'] == 'DEMO-000001'

    demo_sensitive_data = SensitiveData.objects.filter(
        title__startswith=DEMO_MARKER
    )
    assert demo_sensitive_data.count() == 3
    attached_sensitive_data = demo_sensitive_data.get(
        file_attachment__isnull=False
    )
    attached_sensitive_data.file_attachment.refresh_from_db()
    assert (
        attached_sensitive_data.file_attachment.level
        == FileAsset.Level.SENSITIVE
    )
    sensitive_requests = SensitiveAccessRequest.objects.filter(
        project_id__in=demo_project_ids
    )
    assert sensitive_requests.count() == 4
    approved_download = sensitive_requests.get(
        status=SensitiveAccessRequest.Status.APPROVED,
        is_download=True,
    )
    assert approved_download.is_accessible
    assert approved_download.can_download_attachment
    approved_view = sensitive_requests.get(
        status=SensitiveAccessRequest.Status.APPROVED,
        is_download=False,
    )
    assert approved_view.is_accessible
    assert not approved_view.can_download_attachment
    assert approved_view.viewed_at is not None
    assert sensitive_requests.filter(
        status=SensitiveAccessRequest.Status.PENDING
    ).count() == 1
    expired_request = sensitive_requests.get(
        status=SensitiveAccessRequest.Status.EXPIRED
    )
    assert not expired_request.is_accessible

    demo_reports = CustomReport.objects.filter(
        name__startswith=DEMO_MARKER
    )
    demo_schedules = ScheduledReport.objects.filter(report__in=demo_reports)
    demo_executions = ScheduledReportExecution.objects.filter(
        schedule__in=demo_schedules
    )
    assert demo_reports.count() == 3
    assert demo_schedules.count() == 3
    assert demo_executions.count() == 3
    assert set(demo_schedules.values_list('file_format', flat=True)) == {
        ScheduledReport.FileFormat.XLSX,
        ScheduledReport.FileFormat.DOCX,
        ScheduledReport.FileFormat.PDF,
    }
    assert not demo_schedules.exclude(
        last_status=ScheduledReport.RunStatus.SUCCESS
    ).exists()
    for execution in demo_executions.select_related(
        'schedule',
        'schedule__report',
    ):
        assert execution.status == ScheduledReport.RunStatus.SUCCESS
        assert execution.file
        content = Path(execution.file.path).read_bytes()
        assert content
        assert execution.file_size == len(content)
        if execution.file_format == ScheduledReport.FileFormat.XLSX:
            workbook = load_workbook(BytesIO(content), data_only=True)
            parsed_text = '\n'.join(
                str(cell.value or '')
                for row in workbook.active.iter_rows()
                for cell in row
            )
        elif execution.file_format == ScheduledReport.FileFormat.DOCX:
            parsed_text = '\n'.join(
                paragraph.text
                for paragraph in Document(BytesIO(content)).paragraphs
            )
        else:
            parsed_text = '\n'.join(
                page.extract_text() or ''
                for page in PdfReader(BytesIO(content)).pages
            )
        assert 'TEAM-DEMO-' in parsed_text

    demo_announcements = Announcement.objects.filter(
        title__startswith=DEMO_MARKER
    )
    assert demo_announcements.count() == 3
    assert not demo_announcements.exclude(
        status=Announcement.Status.PUBLISHED
    ).exists()
    assert demo_announcements.filter(is_public=True).count() == 1
    assert demo_announcements.filter(is_pinned=True).count() == 1

    demo_imports = ImportTask.objects.filter(
        file_path__contains=str(
            Path('imports') / DEMO_IMPORT_DIRNAME
        )
    )
    assert demo_imports.count() == 8
    assert set(demo_imports.values_list('module', flat=True)) == set(
        ImportTask.Module.values
    )
    assert all(Path(path).exists() for path in demo_imports.values_list(
        'file_path', flat=True
    ))

    demo_contributions = Contribution.objects.filter(
        project_id__in=demo_project_ids
    ).select_related('project', 'user')
    assert demo_contributions.count() == 15
    for contribution in demo_contributions:
        expected_prefix = (
            f'{contribution.user.name}在「{contribution.project.name}」中'
        )
        assert contribution.content.startswith(expected_prefix)
        assert contribution.description == contribution.content

    assert PortalPublication.objects.filter(
        content_type=PortalPublication.ContentType.PROJECT,
        object_id__in=demo_project_ids,
    ).count() == 7
    assert PortalPublication.objects.filter(
        content_type=PortalPublication.ContentType.IP_APPLICATION,
        object_id__in=IntellectualPropertyApplication.objects.filter(
            application_code__startswith=DEMO_IP_PREFIX
        ).values_list('id', flat=True),
    ).count() == 72
    demo_ip = IntellectualPropertyApplication.objects.filter(
        application_code__startswith=DEMO_IP_PREFIX
    ).select_related(
        'related_project',
        'project_reviewer',
        'teacher_confirmer',
        'final_certificate_file',
    )
    assert demo_ip.count() == 72
    patent_types = {
        IntellectualPropertyApplication.IPType.INVENTION_PATENT,
        IntellectualPropertyApplication.IPType.UTILITY_MODEL,
        IntellectualPropertyApplication.IPType.DESIGN_PATENT,
    }
    assert demo_ip.filter(ip_type__in=patent_types).count() == 40
    assert demo_ip.filter(
        ip_type=IntellectualPropertyApplication.IPType.SOFTWARE_COPYRIGHT
    ).count() == 25
    assert demo_ip.filter(
        ip_type=IntellectualPropertyApplication.IPType.NOVELTY_SEARCH
    ).count() == 7
    for application in demo_ip:
        assert application.main_writer_id
        assert application.applicant_executor_id
        assert application.material_manager_id
        assert application.project_reviewer_id == application.related_project.leader_id
        assert application.teacher_confirmer.global_role in {
            User.GlobalRole.TEACHER,
            User.GlobalRole.SYS_ADMIN,
        }
    authorized_ip = demo_ip.get(
        final_certificate_file__isnull=False
    )
    assert authorized_ip.status == (
        IntellectualPropertyApplication.Status.AUTHORIZED
    )
    assert authorized_ip.final_certificate_file.level == FileAsset.Level.INTERNAL
    assert authorized_ip.final_certificate_file.name.endswith(
        '最终授权登记证书.pdf'
    )
    assert Path(authorized_ip.final_certificate_file.file.path).exists()
    certificate_text = '\n'.join(
        page.extract_text() or ''
        for page in PdfReader(
            authorized_ip.final_certificate_file.file.path
        ).pages
    )
    assert '知识产权授权登记证书' in certificate_text
    assert authorized_ip.application_code in certificate_text
    assert authorized_ip.title in certificate_text
    demo_ip_contributors = IPApplicationContributor.objects.filter(
        application__in=demo_ip
    ).select_related('application', 'user')
    assert demo_ip_contributors.count() == 288
    for contributor in demo_ip_contributors:
        assert contributor.contribution_description.startswith(
            f'{contributor.user.name}在「{contributor.application.title}」中'
        )
        assert contributor.responsibility_description.startswith(
            f'{contributor.user.name}负责「{contributor.application.title}」'
        )
    assert not IPApplicationContributor.objects.filter(
        application__in=demo_ip,
        responsibility_description='',
    ).exists()
    assert IPMaterialVersion.objects.filter(application__in=demo_ip).count() == 72
    final_material = IPMaterialVersion.objects.get(
        application=authorized_ip,
        is_final=True,
    )
    assert Path(final_material.file_asset.file.path).exists()
    assert final_material.material_type == IPMaterialVersion.MaterialType.ARCHIVE
    assert (
        final_material.file_asset_id
        != authorized_ip.final_certificate_file_id
    )
    member_publications = PortalPublication.objects.filter(
        content_type=PortalPublication.ContentType.MEMBER,
        custom_summary__startswith=DEMO_MARKER,
    )
    assert member_publications.count() == 17
    assert member_publications.filter(
        is_public=True,
        member_consent=True,
    ).count() == 12
    assert PortalSettings.objects.get(singleton_key='default').team_name == '真实团队门户'

    deterministic_member_skills = set(
        MemberSkill.objects.filter(user__email__in=owned_emails).values_list(
            'user__email',
            'skill__name',
            'proficiency',
        )
    )
    assert 40 <= len(deterministic_member_skills) <= 60
    stale_demo_skill = MemberSkill.objects.create(
        user=User.objects.get(email='admin@demo.com'),
        skill=real_skill,
        proficiency=1,
    )
    real_member_skill = MemberSkill.objects.create(
        user=real_user,
        skill=real_skill,
        proficiency=5,
    )

    # 再运行一次验证精准清理与可重复生成，不得累积本命令数据。
    call_command(
        'seed_demo_data',
        clean=True,
        force=True,
        stdout=StringIO(),
    )
    assert Project.all_objects.filter(
        code__startswith=DEMO_PROJECT_PREFIX
    ).count() == 7
    assert Task.all_objects.filter(
        project__code__startswith=DEMO_PROJECT_PREFIX
    ).count() == 35
    assert FileAsset.objects.filter(
        project__code__startswith=DEMO_PROJECT_PREFIX
    ).count() == 247
    assert ImportTask.objects.filter(
        file_path__contains=str(Path('imports') / DEMO_IMPORT_DIRNAME)
    ).count() == 8
    regenerated_team = Team.objects.get(code=DEMO_TEAM_CODE)
    assert TeamMember.objects.filter(team=regenerated_team).count() == 46
    assert TeamMembershipEvent.objects.filter(
        membership__team=regenerated_team
    ).count() == 53
    regenerated_squads = Team.objects.filter(parent=regenerated_team)
    assert regenerated_squads.count() == 0
    assert all(
        project.teams.filter(pk=regenerated_team.pk).exists()
        for project in Project.all_objects.filter(
            code__startswith=DEMO_PROJECT_PREFIX
        )
    )
    assert FinanceReceipt.objects.filter(
        expense__project__code__startswith=DEMO_PROJECT_PREFIX
    ).count() == expected_receipt_count
    assert UserPreference.objects.filter(user__email__in=owned_emails).count() == 47
    assert Competition.objects.filter(
        project__code__startswith=DEMO_PROJECT_PREFIX
    ).count() == 109
    assert CompetitionEvent.objects.filter(
        organization=regenerated_team
    ).count() == 33
    assert IntellectualPropertyApplication.objects.filter(
        application_code__startswith=DEMO_IP_PREFIX
    ).count() == 72
    assert SensitiveAccessRequest.objects.filter(
        project__code__startswith=DEMO_PROJECT_PREFIX
    ).count() == 4
    assert CustomReport.objects.filter(
        name__startswith=DEMO_MARKER
    ).count() == 3
    assert ScheduledReport.objects.filter(
        report__name__startswith=DEMO_MARKER
    ).count() == 3
    regenerated_executions = ScheduledReportExecution.objects.filter(
        schedule__report__name__startswith=DEMO_MARKER
    )
    assert regenerated_executions.count() == 3
    assert all(
        Path(execution.file.path).exists()
        for execution in regenerated_executions
    )
    assert Announcement.objects.filter(
        title__startswith=DEMO_MARKER
    ).count() == 3
    assert not MemberSkill.objects.filter(pk=stale_demo_skill.pk).exists()
    assert set(
        MemberSkill.objects.filter(user__email__in=owned_emails).values_list(
            'user__email',
            'skill__name',
            'proficiency',
        )
    ) == deterministic_member_skills

    # 真实数据和全局门户配置必须完整保留。
    assert User.objects.filter(pk=real_user.pk).exists()
    assert Team.objects.filter(pk=real_team.pk).exists()
    assert TeamMember.objects.filter(pk=real_team_member.pk).exists()
    assert TeamMembershipEvent.objects.filter(pk=real_team_event.pk).exists()
    assert Project.all_objects.filter(pk=real_project.pk).exists()
    assert Task.all_objects.filter(pk=real_task.pk).exists()
    assert FinanceIncome.objects.filter(pk=real_income.pk).exists()
    assert FileAsset.objects.filter(pk=real_file.pk).exists()
    assert Path(real_file.file.path).exists()
    assert ImportTask.objects.filter(pk=real_import.pk).exists()
    assert SkillTag.objects.filter(pk=real_skill.pk).exists()
    assert MemberSkill.objects.filter(pk=real_member_skill.pk).exists()
    assert PortalSettings.objects.filter(pk=portal_settings.pk).exists()
    assert PortalPublication.objects.filter(pk=real_publication.pk).exists()
    assert UserLifecycleEvent.objects.filter(pk=real_lifecycle.pk).exists()
    assert CustomReport.objects.filter(pk=real_report.pk).exists()
    assert ScheduledReport.objects.filter(pk=real_schedule.pk).exists()
    assert ScheduledReportExecution.objects.filter(pk=real_execution.pk).exists()
    assert Path(real_execution.file.path).exists()
    assert Announcement.objects.filter(pk=real_announcement.pk).exists()


@pytest.mark.django_db(transaction=True)
def test_seed_commands_share_one_complete_dataset_in_both_orders(settings, tmp_path):
    settings.MEDIA_ROOT = tmp_path / 'media'

    def assert_complete_dataset():
        projects = Project.all_objects.filter(code__startswith=DEMO_PROJECT_PREFIX)
        assert set(
            User.objects.filter(email__endswith='@demo.com').values_list(
                'email', flat=True
            )
        ) == set(DEMO_ACCOUNT_EMAILS)
        assert list(
            User.objects.filter(
                global_role=User.GlobalRole.TEACHER,
                is_active=True,
            ).values_list('email', flat=True)
        ) == ['teacher1@demo.com']
        read_only_teacher_emails = {
            'teacher2@demo.com',
            'teacher3@demo.com',
            'teacher4@demo.com',
        }
        assert set(
            User.objects.filter(
                email__in=read_only_teacher_emails,
                global_role=User.GlobalRole.MEMBER,
            ).values_list('email', flat=True)
        ) == read_only_teacher_emails
        demo_team = Team.objects.get(code=DEMO_TEAM_CODE)
        assert set(
            TeamMember.objects.filter(
                team=demo_team,
                role=TeamMember.Role.TEACHER,
                user__email__in={
                    'teacher1@demo.com',
                    *read_only_teacher_emails,
                },
            ).values_list('user__email', flat=True)
        ) == {
            'teacher1@demo.com',
            *read_only_teacher_emails,
        }
        assert projects.count() == 7
        assert Task.all_objects.filter(project__in=projects).count() == 35
        assert Competition.objects.filter(project__in=projects).count() == 109
        assert IntellectualPropertyApplication.objects.filter(
            application_code__startswith=DEMO_IP_PREFIX
        ).count() == 72
        demo_reports = CustomReport.objects.filter(name__startswith=DEMO_MARKER)
        assert demo_reports.count() == 3
        assert ScheduledReport.objects.filter(report__in=demo_reports).count() == 3
        assert not CustomReport.objects.filter(
            name__in=LEGACY_COMPETITION_REPORT_NAMES
        ).exists()

    legacy_users = [
        User.objects.create_user(
            email=email,
            username=email.split('@', 1)[0],
            password='legacy-demo-password',
            name='旧比赛演示账号',
        )
        for email in LEGACY_COMPETITION_ACCOUNT_EMAILS
    ]
    for index, name in enumerate(LEGACY_COMPETITION_REPORT_NAMES):
        report = CustomReport.objects.create(
            name=name,
            report_type='project',
            created_by=legacy_users[index % len(legacy_users)],
            is_scheduled=True,
        )
        ScheduledReport.objects.create(
            report=report,
            created_by=report.created_by,
        )

    call_command(
        'seed_competition_demo',
        clean=True,
        force=True,
        stdout=StringIO(),
    )
    assert_complete_dataset()

    call_command(
        'seed_demo_data',
        clean=True,
        force=True,
        stdout=StringIO(),
    )
    assert_complete_dataset()

    call_command(
        'seed_competition_demo',
        clean=True,
        force=True,
        stdout=StringIO(),
    )
    assert_complete_dataset()
