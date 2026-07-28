from io import BytesIO

import pytest
from django.core.files.uploadedfile import SimpleUploadedFile

from apps.files.models import FileAsset
from apps.files.share_models import FileShareLink


def extract_data(response):
    body = response.json()
    if isinstance(body, dict) and 'code' in body:
        return body.get('data', body)
    return body


def extract_results(response):
    data = extract_data(response)
    if isinstance(data, dict) and 'results' in data:
        return data['results']
    return data


@pytest.mark.api
@pytest.mark.django_db
class TestFileFolderClosure:
    def test_folder_crud_and_file_move_are_project_scoped(
        self,
        teacher_client,
        make_project,
        make_file,
    ):
        project = make_project(name='文件项目')
        other_project = make_project(name='其他项目')

        root_response = teacher_client.post(
            '/api/v1/files/folders/',
            {'project': project.id, 'name': '交付材料'},
            format='json',
        )
        assert root_response.status_code == 201, root_response.json()
        root = extract_data(root_response)

        duplicate = teacher_client.post(
            '/api/v1/files/folders/',
            {'project': project.id, 'name': '交付材料'},
            format='json',
        )
        assert duplicate.status_code == 400

        child_response = teacher_client.post(
            '/api/v1/files/folders/',
            {'project': project.id, 'parent': root['id'], 'name': '终稿'},
            format='json',
        )
        assert child_response.status_code == 201, child_response.json()
        child = extract_data(child_response)
        assert child['path'] == '交付材料 / 终稿'

        upload_response = teacher_client.post(
            '/api/v1/files/',
            {
                'project': project.id,
                'folder': child['id'],
                'level': 'internal',
                'file': SimpleUploadedFile(
                    '目录关联.txt',
                    b'folder-linked-content',
                    content_type='text/plain',
                ),
            },
            format='multipart',
        )
        assert upload_response.status_code == 201, upload_response.json()
        uploaded = extract_data(upload_response)
        assert uploaded['folder'] == child['id']
        assert uploaded['folder_name'] == '终稿'
        assert FileAsset.objects.get(pk=uploaded['id']).folder_id == child['id']

        cross_project = teacher_client.post(
            '/api/v1/files/folders/',
            {
                'project': other_project.id,
                'parent': root['id'],
                'name': '越界目录',
            },
            format='json',
        )
        assert cross_project.status_code == 400

        file_asset = make_file(project=project)
        moved = teacher_client.post(
            f'/api/v1/files/{file_asset.id}/move/',
            {'folder': child['id']},
            format='json',
        )
        assert moved.status_code == 200, moved.json()
        assert extract_data(moved)['folder'] == child['id']

        filtered = teacher_client.get(
            f'/api/v1/files/?project={project.id}&folder={child["id"]}',
        )
        assert filtered.status_code == 200
        assert {
            item['id'] for item in extract_results(filtered)
        } == {uploaded['id'], file_asset.id}

        moved_to_root = teacher_client.post(
            f'/api/v1/files/{file_asset.id}/move/',
            {'folder': None},
            format='json',
        )
        assert moved_to_root.status_code == 200
        assert extract_data(moved_to_root)['folder'] is None

        renamed = teacher_client.patch(
            f'/api/v1/files/folders/{root["id"]}/',
            {'name': '验收材料'},
            format='json',
        )
        assert renamed.status_code == 200, renamed.json()
        assert extract_data(renamed)['name'] == '验收材料'

    def test_folder_cannot_be_moved_below_its_child(
        self,
        teacher_client,
        make_project,
    ):
        project = make_project()
        root = extract_data(teacher_client.post(
            '/api/v1/files/folders/',
            {'project': project.id, 'name': '根目录'},
            format='json',
        ))
        child = extract_data(teacher_client.post(
            '/api/v1/files/folders/',
            {'project': project.id, 'parent': root['id'], 'name': '子目录'},
            format='json',
        ))
        response = teacher_client.patch(
            f'/api/v1/files/folders/{root["id"]}/',
            {'parent': child['id']},
            format='json',
        )
        assert response.status_code == 400


@pytest.mark.api
@pytest.mark.django_db
class TestFileRecycleClosure:
    def test_delete_restore_and_permanent_delete(
        self,
        teacher_client,
        admin_client,
        make_file,
    ):
        file_asset = make_file(uploader=teacher_client.user)
        share = FileShareLink.objects.create(
            file=file_asset,
            created_by=teacher_client.user,
            token=FileShareLink.generate_token(),
        )

        deleted = teacher_client.delete(f'/api/v1/files/{file_asset.id}/')
        assert deleted.status_code == 200, deleted.json()
        recycled = FileAsset.all_objects.get(pk=file_asset.pk)
        assert recycled.is_deleted is True
        assert recycled.deleted_by_id == teacher_client.user.id
        share.refresh_from_db()
        assert share.is_active is False

        recycle_list = teacher_client.get('/api/v1/recycle-bin/?type=file')
        assert recycle_list.status_code == 200, recycle_list.json()
        assert file_asset.id in [row['id'] for row in extract_data(recycle_list)]

        restored = teacher_client.post(
            '/api/v1/recycle-bin/',
            {'type': 'file', 'id': file_asset.id},
            format='json',
        )
        assert restored.status_code == 200, restored.json()
        assert FileAsset.objects.filter(pk=file_asset.pk).exists()
        share.refresh_from_db()
        assert share.is_active is False

        teacher_client.delete(f'/api/v1/files/{file_asset.id}/')
        permanently_deleted = admin_client.delete(
            f'/api/v1/recycle-bin/?type=file&id={file_asset.id}',
        )
        assert permanently_deleted.status_code == 200, permanently_deleted.json()
        assert not FileAsset.all_objects.filter(pk=file_asset.pk).exists()


def _docx_bytes():
    from docx import Document

    output = BytesIO()
    document = Document()
    document.add_heading('项目周报', level=1)
    document.add_paragraph('本周已经完成接口联调。')
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = '负责人'
    table.cell(0, 1).text = '张老师'
    document.save(output)
    return output.getvalue()


def _xlsx_bytes():
    from openpyxl import Workbook

    output = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = '进度'
    sheet.append(['任务', '状态'])
    sheet.append(['接口联调', '完成'])
    workbook.save(output)
    return output.getvalue()


def _pptx_bytes():
    from pptx import Presentation

    output = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = '项目答辩'
    slide.placeholders[1].text = '核心指标已经达标'
    presentation.save(output)
    return output.getvalue()


@pytest.mark.api
@pytest.mark.django_db
@pytest.mark.parametrize(
    ('name', 'content_type', 'builder', 'expected_type', 'expected_text'),
    [
        (
            '周报.docx',
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            _docx_bytes,
            'docx',
            '本周已经完成接口联调。',
        ),
        (
            '进度.xlsx',
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            _xlsx_bytes,
            'xlsx',
            '接口联调',
        ),
        (
            '答辩.pptx',
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            _pptx_bytes,
            'pptx',
            '核心指标已经达标',
        ),
    ],
)
def test_office_preview_extracts_bounded_plain_content(
    teacher_client,
    make_project,
    name,
    content_type,
    builder,
    expected_type,
    expected_text,
):
    payload = builder()
    file_asset = FileAsset.objects.create(
        project=make_project(),
        name=name,
        file=SimpleUploadedFile(name, payload, content_type=content_type),
        size=len(payload),
        content_type=content_type,
        uploader=teacher_client.user,
    )

    response = teacher_client.get(
        f'/api/v1/files/{file_asset.id}/office-preview/',
    )
    assert response.status_code == 200, response.json()
    data = extract_data(response)
    assert data['type'] == expected_type
    assert expected_text in str(data['sections'])
    assert data['limits']['source_bytes'] == 10 * 1024 * 1024


@pytest.mark.api
@pytest.mark.django_db
def test_office_preview_rejects_macro_enabled_and_sensitive_files(
    teacher_client,
    make_project,
):
    payload = _docx_bytes()
    macro_file = FileAsset.objects.create(
        project=make_project(),
        name='含宏文档.docm',
        file=SimpleUploadedFile('含宏文档.docm', payload),
        size=len(payload),
        uploader=teacher_client.user,
    )
    rejected = teacher_client.get(
        f'/api/v1/files/{macro_file.id}/office-preview/',
    )
    assert rejected.status_code == 400

    sensitive_file = FileAsset.objects.create(
        project=make_project(),
        name='敏感周报.docx',
        file=SimpleUploadedFile('敏感周报.docx', payload),
        size=len(payload),
        level=FileAsset.Level.SENSITIVE,
        uploader=teacher_client.user,
    )
    forbidden = teacher_client.get(
        f'/api/v1/files/{sensitive_file.id}/office-preview/',
    )
    assert forbidden.status_code == 403
